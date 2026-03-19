from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ══════════════════════════════════════════════════════════════
# COLOR PALETTE
# ══════════════════════════════════════════════════════════════
C_PRIMARY = RGBColor(59, 130, 246)
C_SECONDARY = RGBColor(139, 92, 246)
C_ACCENT1 = RGBColor(16, 185, 129)
C_ACCENT2 = RGBColor(245, 158, 11)
C_ACCENT3 = RGBColor(239, 68, 68)
C_ACCENT4 = RGBColor(6, 182, 212)
C_DARK = RGBColor(30, 41, 59)
C_MID = RGBColor(71, 85, 105)
C_WHITE = RGBColor(255, 255, 255)

def add_bg(slide, color=C_WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=C_DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_down_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body, icon_color, title_color=C_DARK):
    add_rounded_rect(slide, left, top, width, height, C_WHITE)
    add_rect(slide, left, top, width, Inches(0.06), icon_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.35),
             title, size=13, color=title_color, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
             body, size=10, color=C_MID)

def add_slide_number(slide, num, total=10):
    add_text(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
             f"{num}/{total}", size=10, color=C_MID, align=PP_ALIGN.RIGHT)

def add_header_bar(slide, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), C_PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
             title, size=28, color=C_WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4),
                 subtitle, size=14, color=RGBColor(191, 219, 254))
    add_circle(slide, Inches(11.5), Inches(-0.3), Inches(0.8), RGBColor(96, 165, 250))
    add_circle(slide, Inches(12.2), Inches(0.3), Inches(0.6), C_SECONDARY)

# ══════════════════════════════════════════════════════════════
# SLIDE 1: DEPLOYMENT (AWS + Jenkins CI/CD)
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, C_WHITE)
add_header_bar(slide1, "Deployment Pipeline", "AWS Cloud Deployment with Jenkins CI/CD for Monolithic Version")
add_slide_number(slide1, 1)

# CI/CD Flow — horizontal pipeline
stages = [
    ("Developer\nPushes Code", "Git commit & push\nto GitHub repo", RGBColor(99, 102, 241)),
    ("GitHub\nWebhook", "Triggers Jenkins\nbuild pipeline", C_DARK),
    ("Jenkins\nCI Pipeline", "Maven build, run\nunit tests, package", C_ACCENT2),
    ("Docker\nBuild", "Dockerfile builds\napp container image", C_ACCENT4),
    ("Push to\nDocker Hub / ECR", "Image pushed to\ncontainer registry", C_SECONDARY),
    ("AWS EC2\nDeploy", "Jenkins deploys to\nEC2 via SSH/ECS", C_ACCENT1),
]

for i, (title, desc, color) in enumerate(stages):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(1.4)
    add_rounded_rect(slide1, x, y, Inches(1.95), Inches(1.6), color)
    add_text(slide1, x + Inches(0.1), y + Inches(0.15), Inches(1.75), Inches(0.5),
             title, size=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, x + Inches(0.1), y + Inches(0.7), Inches(1.75), Inches(0.7),
             desc, size=10, color=RGBColor(220, 230, 255), align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_arrow(slide1, x + Inches(1.95), y + Inches(0.6), Inches(0.2), Inches(0.15), RGBColor(203, 213, 225))

# AWS Architecture below
add_text(slide1, Inches(0.5), Inches(3.3), Inches(12), Inches(0.4),
         "AWS Cloud Architecture", size=18, color=C_DARK, bold=True)

aws_components = [
    ("Route 53 / DNS", "Domain routing to\nload balancer", RGBColor(99, 102, 241)),
    ("Elastic Load\nBalancer (ALB)", "Distributes traffic\nacross instances", C_ACCENT2),
    ("EC2 Instance", "Runs Spring Boot\napp in Docker", C_ACCENT1),
    ("RDS MySQL", "Managed MySQL\ndatabase service", C_PRIMARY),
    ("S3 Bucket", "Static assets &\nmedia file storage", C_ACCENT4),
]

for i, (title, desc, color) in enumerate(aws_components):
    x = Inches(0.3) + Inches(i * 2.6)
    y = Inches(3.8)
    add_rounded_rect(slide1, x, y, Inches(2.3), Inches(1.2), C_WHITE)
    add_rect(slide1, x, y, Inches(2.3), Inches(0.06), color)
    add_text(slide1, x + Inches(0.15), y + Inches(0.15), Inches(2.0), Inches(0.4),
             title, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, x + Inches(0.15), y + Inches(0.55), Inches(2.0), Inches(0.5),
             desc, size=9, color=C_MID, align=PP_ALIGN.CENTER)
    if i < len(aws_components) - 1:
        add_arrow(slide1, x + Inches(2.3), y + Inches(0.45), Inches(0.3), Inches(0.15), RGBColor(203, 213, 225))

# Jenkins details
add_text(slide1, Inches(0.5), Inches(5.3), Inches(12), Inches(0.3),
         "Jenkins Pipeline Stages (Jenkinsfile)", size=14, color=C_DARK, bold=True)

jenkins_cards = [
    ("Checkout", "git checkout from\nGitHub repository", RGBColor(99, 102, 241)),
    ("Build", "mvn clean package\n-DskipTests", C_ACCENT2),
    ("Test", "mvn test — runs\nunit & integration tests", C_ACCENT1),
    ("Docker Build", "docker build -t\nrevconnect:latest .", C_ACCENT4),
    ("Push Image", "docker push to\nDocker Hub / AWS ECR", C_SECONDARY),
    ("Deploy", "ssh to EC2, pull image\ndocker-compose up -d", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(jenkins_cards):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(5.7)
    add_rounded_rect(slide1, x, y, Inches(1.95), Inches(1.1), C_WHITE)
    add_rect(slide1, x, y, Inches(1.95), Inches(0.05), color)
    add_circle(slide1, x + Inches(0.1), y + Inches(0.15), Inches(0.3), color)
    add_text(slide1, x + Inches(0.1), y + Inches(0.17), Inches(0.3), Inches(0.3),
             str(i + 1), size=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, x + Inches(0.5), y + Inches(0.12), Inches(1.3), Inches(0.25),
             title, size=10, color=color, bold=True)
    add_text(slide1, x + Inches(0.5), y + Inches(0.4), Inches(1.3), Inches(0.6),
             desc, size=9, color=C_MID)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: INTRODUCTION TO REVCONNECT
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, C_WHITE)

# Left panel
add_rect(slide2, Inches(0), Inches(0), Inches(5.5), prs.slide_height, RGBColor(238, 242, 255))
add_rect(slide2, Inches(0), Inches(0), Inches(0.08), prs.slide_height, C_PRIMARY)

add_circle(slide2, Inches(10), Inches(0.5), Inches(2.5), RGBColor(238, 242, 255))
add_circle(slide2, Inches(11), Inches(5), Inches(2), RGBColor(245, 243, 255))
add_circle(slide2, Inches(4.5), Inches(5.5), Inches(1.5), RGBColor(236, 253, 245))

add_text(slide2, Inches(0.8), Inches(0.8), Inches(4), Inches(0.5),
         "RevConnect", size=42, color=C_PRIMARY, bold=True)
add_text(slide2, Inches(0.8), Inches(1.4), Inches(4), Inches(0.5),
         "Social Media Platform", size=22, color=C_SECONDARY, bold=True)

add_text(slide2, Inches(0.8), Inches(2.2), Inches(4.2), Inches(0.3),
         "ABSTRACT", size=14, color=C_PRIMARY, bold=True)

abstract_text = (
    "RevConnect is a full-stack professional social media platform built using a microservices architecture "
    "with Angular 19 frontend and Spring Boot 3 backend. The system comprises 7 independent microservices "
    "(User, Post, Interaction, Connection, Notification, Feed, API Gateway) each with its own MySQL database, "
    "communicating via REST APIs and Eureka service discovery.\n\n"
    "Key features include user registration with email verification, JWT-based authentication, real-time "
    "personalized feeds, post creation with media uploads (images/videos), social interactions (likes, comments, "
    "shares), connection management (follow/unfollow), direct messaging with WebSocket support, push/email "
    "notifications with granular preferences, user analytics dashboards, stories, bookmarks, explore/search, "
    "and business tools (CTA buttons, paid partnerships, product tagging, scheduled posts).\n\n"
    "The platform supports three user types: Standard, Business, and Creator — each with tailored features. "
    "Privacy controls allow users to set profiles as Public or Private. The architecture ensures scalability, "
    "fault tolerance via circuit breakers, and clean separation of concerns."
)
add_text(slide2, Inches(0.8), Inches(2.6), Inches(4.2), Inches(4.5),
         abstract_text, size=11, color=C_MID)

# Tech stack cards
add_text(slide2, Inches(6), Inches(0.8), Inches(6), Inches(0.4),
         "Technology Stack", size=20, color=C_DARK, bold=True)

techs = [
    ("Frontend", "Angular 19, TypeScript, RxJS, TailwindCSS, Font Awesome", C_PRIMARY),
    ("Backend", "Spring Boot 3, Java 17, Spring Cloud, Eureka, Feign, JWT", C_SECONDARY),
    ("Database", "MySQL 8 (5 Docker containers), JPA/Hibernate, Spring Data", C_ACCENT1),
    ("DevOps", "Docker Compose, API Gateway, Circuit Breakers, Eureka Discovery", C_ACCENT2),
    ("Features", "Real-time Feed, Analytics, Stories, Messages, Notifications, Business Tools", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(techs):
    y = Inches(1.4) + Inches(i * 1.1)
    add_card(slide2, Inches(6), y, Inches(6.8), Inches(0.9), title, desc, color)

add_slide_number(slide2, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, C_WHITE)
add_header_bar(slide3, "System Architecture", "Microservices Architecture with Service Discovery & API Gateway")
add_slide_number(slide3, 3)

# Client
add_rounded_rect(slide3, Inches(5.2), Inches(1.3), Inches(3), Inches(0.7), RGBColor(219, 234, 254))
add_text(slide3, Inches(5.2), Inches(1.35), Inches(3), Inches(0.35),
         "Angular 19 Frontend", size=14, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, Inches(5.2), Inches(1.65), Inches(3), Inches(0.3),
         "localhost:4200", size=10, color=C_MID, align=PP_ALIGN.CENTER)

add_arrow(slide3, Inches(6.4), Inches(2.1), Inches(0.6), Inches(0.3), C_MID)

# Gateway
add_rounded_rect(slide3, Inches(4.5), Inches(2.5), Inches(4.4), Inches(0.8), RGBColor(254, 243, 199))
add_text(slide3, Inches(4.5), Inches(2.55), Inches(4.4), Inches(0.4),
         "API Gateway (Spring Cloud Gateway)", size=14, color=C_ACCENT2, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, Inches(4.5), Inches(2.9), Inches(4.4), Inches(0.3),
         "JWT Auth | Route | Load Balance | Circuit Breaker — :8080", size=10, color=C_MID, align=PP_ALIGN.CENTER)

# Eureka
add_rounded_rect(slide3, Inches(10), Inches(2.5), Inches(2.8), Inches(0.8), RGBColor(237, 233, 254))
add_text(slide3, Inches(10), Inches(2.55), Inches(2.8), Inches(0.4),
         "Eureka Server", size=14, color=C_SECONDARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, Inches(10), Inches(2.9), Inches(2.8), Inches(0.3),
         "Service Discovery — :8761", size=10, color=C_MID, align=PP_ALIGN.CENTER)

# Services
services = [
    ("User\nService", ":8081", C_PRIMARY, "Auth, Profile\nSettings, Search"),
    ("Post\nService", ":8082", C_ACCENT1, "CRUD, Schedule\nMedia, CTA"),
    ("Feed\nService", ":8083", C_ACCENT4, "Personalized\nFeed Aggregation"),
    ("Interaction\nService", ":8084", C_ACCENT2, "Likes, Comments\nShares, Views"),
    ("Connection\nService", ":8085", C_SECONDARY, "Follow/Unfollow\nStats, Network"),
    ("Notification\nService", ":8086", C_ACCENT3, "Push, Email\nPreferences"),
]
for i, (name, port, color, desc) in enumerate(services):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(4.0)
    add_rounded_rect(slide3, x, y, Inches(1.95), Inches(1.8), C_WHITE)
    add_rect(slide3, x, y, Inches(1.95), Inches(0.08), color)
    add_text(slide3, x, y + Inches(0.15), Inches(1.95), Inches(0.5),
             name, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, x, y + Inches(0.65), Inches(1.95), Inches(0.25),
             port, size=9, color=C_MID, align=PP_ALIGN.CENTER)
    add_text(slide3, x, y + Inches(0.9), Inches(1.95), Inches(0.7),
             desc, size=9, color=C_MID, align=PP_ALIGN.CENTER)

# DBs
dbs = [
    ("user_db", ":3307", C_PRIMARY),
    ("post_db", ":3308", C_ACCENT1),
    ("—", "", C_ACCENT4),
    ("interaction_db", ":3309", C_ACCENT2),
    ("connection_db", ":3310", C_SECONDARY),
    ("notification_db", ":3311", C_ACCENT3),
]
for i, (name, port, color) in enumerate(dbs):
    if name == "—": continue
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(6.1)
    add_rounded_rect(slide3, x, y, Inches(1.95), Inches(0.7), RGBColor(241, 245, 249))
    add_text(slide3, x, y + Inches(0.05), Inches(1.95), Inches(0.35),
             f"MySQL {name}", size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, x, y + Inches(0.35), Inches(1.95), Inches(0.25),
             f"Docker {port}", size=8, color=C_MID, align=PP_ALIGN.CENTER)

add_text(slide3, Inches(0.3), Inches(6.9), Inches(12), Inches(0.3),
         "Each microservice has its own isolated MySQL database running in Docker containers — Database per Service pattern",
         size=10, color=C_MID, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: LOGIN, REGISTER, FORGOT PASSWORD
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, C_WHITE)
add_header_bar(slide4, "Authentication Flow", "Login, Register & Forgot Password")
add_slide_number(slide4, 4)

# LOGIN
add_text(slide4, Inches(0.5), Inches(1.3), Inches(5), Inches(0.4),
         "Login Flow", size=18, color=C_PRIMARY, bold=True)
login_steps = [
    ("1", "User enters\nemail & password", C_PRIMARY),
    ("2", "Angular sends\nPOST /api/auth/login", C_ACCENT4),
    ("3", "Gateway routes\nto User Service", C_ACCENT2),
    ("4", "BCrypt verifies\npassword hash", C_ACCENT3),
    ("5", "JWT token\ngenerated & returned", C_ACCENT1),
]
for i, (num, text, color) in enumerate(login_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_circle(slide4, x, Inches(1.9), Inches(0.5), color)
    add_text(slide4, x, Inches(1.93), Inches(0.5), Inches(0.5),
             num, size=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide4, x - Inches(0.3), Inches(2.5), Inches(1.1), Inches(0.6),
             text, size=10, color=C_DARK, align=PP_ALIGN.CENTER)
    if i < len(login_steps) - 1:
        add_arrow(slide4, x + Inches(0.6), Inches(2.05), Inches(1.6), Inches(0.2), RGBColor(203, 213, 225))

add_text(slide4, Inches(0.5), Inches(3.2), Inches(12), Inches(0.5),
         "Dependencies: BCryptPasswordEncoder, JwtTokenProvider (HS256, 256-bit secret), Spring Security FilterChain",
         size=10, color=C_MID)

# REGISTER
add_text(slide4, Inches(0.5), Inches(3.8), Inches(5), Inches(0.4),
         "Register Flow", size=18, color=C_ACCENT1, bold=True)
reg_steps = [
    ("1", "User fills form:\nname, email, password", C_ACCENT1),
    ("2", "POST /api/auth/\nregister", C_ACCENT4),
    ("3", "Duplicate check\n+ BCrypt hash", C_ACCENT2),
    ("4", "Verification email\nsent (code in DB)", C_ACCENT3),
    ("5", "User clicks link\naccount activated", C_PRIMARY),
]
for i, (num, text, color) in enumerate(reg_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_circle(slide4, x, Inches(4.4), Inches(0.5), color)
    add_text(slide4, x, Inches(4.43), Inches(0.5), Inches(0.5),
             num, size=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide4, x - Inches(0.3), Inches(5.0), Inches(1.1), Inches(0.6),
             text, size=10, color=C_DARK, align=PP_ALIGN.CENTER)
    if i < len(reg_steps) - 1:
        add_arrow(slide4, x + Inches(0.6), Inches(4.55), Inches(1.6), Inches(0.2), RGBColor(203, 213, 225))

# FORGOT PASSWORD
add_text(slide4, Inches(0.5), Inches(5.8), Inches(5), Inches(0.4),
         "Forgot Password Flow", size=18, color=C_ACCENT3, bold=True)
add_text(slide4, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
         "User enters email \u2192 POST /api/auth/forgot-password \u2192 User Service generates reset token (UUID) \u2192 "
         "Email sent with reset link \u2192 User clicks link with token \u2192 Enters new password \u2192 "
         "POST /api/auth/reset-password with token \u2192 Password re-hashed with BCrypt \u2192 Success",
         size=11, color=C_MID)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: HOME PAGE, FEED, POSTS
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, C_WHITE)
add_header_bar(slide5, "Home Page, Feed & Post Creation", "Personalized feed, scheduling, CTA & business tools")
add_slide_number(slide5, 5)

# Feed flow
flow_items = [
    ("User opens /feed", "FeedPage component\ninitializes", C_PRIMARY),
    ("Load profile", "GET /api/users/me\nvia API Gateway + JWT", C_ACCENT4),
    ("Fetch feed", "GET /api/posts/feed/\npersonalized", C_ACCENT1),
    ("Feed Service", "Gets following IDs from\nConnection Service (Feign)", C_ACCENT2),
    ("Post Service", "Queries posts by\nfollowing user IDs", C_SECONDARY),
    ("Enrich posts", "Fetches author info\nfrom User Service", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(flow_items):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(1.4)
    add_rounded_rect(slide5, x, y, Inches(1.95), Inches(1.3), C_WHITE)
    add_rect(slide5, x, y, Inches(1.95), Inches(0.06), color)
    add_text(slide5, x + Inches(0.1), y + Inches(0.12), Inches(1.75), Inches(0.3),
             title, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide5, x + Inches(0.1), y + Inches(0.45), Inches(1.75), Inches(0.7),
             desc, size=10, color=C_MID, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_arrow(slide5, x + Inches(1.95), y + Inches(0.5), Inches(0.2), Inches(0.15), RGBColor(203, 213, 225))

# Post creation & business tools
add_text(slide5, Inches(0.5), Inches(3.0), Inches(5), Inches(0.3),
         "Post Types & Business Tools", size=16, color=C_DARK, bold=True)

cards_data = [
    ("Normal Post", "Text + media upload.\nPOST /api/posts.\nTypes: TEXT, IMAGE, VIDEO", C_PRIMARY),
    ("Scheduled Post", "Select date/time + CTA.\nPOST /api/posts/schedule.\nAuto-published by scheduler.", C_ACCENT2),
    ("Business CTA", "CTA buttons: Shop Now, Learn More,\nSign Up, Download, Book Now.\nCustom URL link.", C_ACCENT1),
    ("Creator Tools", "Paid Partnership disclosure +\nProduct tagging (comma-sep).\nPartner: Nike, Apple, etc.", C_SECONDARY),
    ("Stories", "24h ephemeral content.\nHorizontal scroll at top.\nImage/video stories.", C_ACCENT3),
    ("Media Upload", "Image & video files.\nPOST /api/media/upload.\nServed from /uploads/ path.", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(cards_data):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(3.4) + Inches(row * 1.8)
    add_card(slide5, x, y, Inches(3.8), Inches(1.6), title, desc, color)

# ══════════════════════════════════════════════════════════════
# SLIDE 6: LIKES, COMMENTS, SHARES, MESSAGES
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, C_WHITE)
add_header_bar(slide6, "Social Interactions & Messaging", "Likes, Comments, Shares, Views & Direct Messages")
add_slide_number(slide6, 6)

interactions = [
    ("Like Flow", "User clicks \u2764 \u2192 POST /api/posts/{id}/like \u2192 Interaction Service checks unique constraint "
     "(user_id + post_id) \u2192 If not liked: saves Like entity + sends notification to post author \u2192 "
     "Unlike: DELETE removes record. Prevents duplicate likes.", C_ACCENT3),
    ("Comment Flow", "User types comment \u2192 POST /api/posts/{id}/comments with {content, parentId?} \u2192 "
     "Supports nested replies via parentId \u2192 Notification sent to author \u2192 "
     "Paginated fetch: GET /api/posts/{id}/comments?page=0&size=10", C_PRIMARY),
    ("Share Flow", "User clicks Share \u2192 POST /api/posts/{id}/share with {comment} \u2192 "
     "Creates Share entity (REPOST type) \u2192 Notification to original author \u2192 "
     "Share count visible in post cards and analytics.", C_ACCENT1),
    ("View Tracking", "Feed loads \u2192 POST /api/posts/{id}/view for each post \u2192 "
     "PostView entity with unique constraint (post_id + user_id) \u2192 "
     "No duplicate counting on refresh. Used in analytics dashboard.", C_ACCENT4),
    ("Direct Messages", "Real-time via WebSocket (STOMP). Types: TEXT, IMAGE, VOICE. "
     "Conversation list with unread count. Read receipts. Typing indicators. "
     "GET /api/messages/conversations, POST /api/messages.", C_SECONDARY),
    ("Bookmarks", "Save posts for later: POST /api/interactions/bookmarks/{postId}. "
     "Toggle bookmark on/off. GET /api/interactions/bookmarks returns saved posts. "
     "Accessible from sidebar.", C_ACCENT2),
]
for i, (title, desc, color) in enumerate(interactions):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.3) + Inches(row * 1.9)
    add_card(slide6, x, y, Inches(5.9), Inches(1.7), title, desc, color)

# ══════════════════════════════════════════════════════════════
# SLIDE 7: EXPLORE, CONNECTIONS, PROFILE, ANALYTICS
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, C_WHITE)
add_header_bar(slide7, "Explore, Connections, Profile & Analytics", "Search, follow/unfollow, profile management & data insights")
add_slide_number(slide7, 7)

features = [
    ("Explore / Search",
     "User types query \u2192 GET /api/search/all?query=... (debounced) \u2192 "
     "SearchController calls UserService + PostService \u2192 "
     "PRIVATE accounts excluded from all results \u2192 "
     "Results shown in All / Users / Posts tabs.", C_PRIMARY),
    ("Follow / Unfollow",
     "Click Follow \u2192 POST /api/users/{id}/follow \u2192 "
     "Connection Service creates record \u2192 FOLLOW notification sent \u2192 "
     "Follower/Following counts updated. Unfollow: DELETE removes record.",  C_ACCENT1),
    ("Profile Page",
     "Cover photo, avatar, name, bio, location, website, user type badge. "
     "Edit via PUT /api/users/me. User posts tab with paginated posts. "
     "Connection stats: followers & following counts. Bookmarks tab.", C_ACCENT2),
    ("Analytics Dashboard",
     "Overview: totalViews, totalLikes, totalComments, totalShares, totalFollowers. "
     "Per-post performance breakdown. Real follower demographics by account type. "
     "Engagement rate calculation. Data aggregated from 3 microservices via Feign.", C_SECONDARY),
    ("Notifications",
     "Granular per-type preferences: likes, comments, shares, followers. "
     "Types: LIKE, COMMENT, SHARE, FOLLOW, MENTION, SYSTEM. "
     "Preferences checked before sending. Stored in UserSettings entity.", C_ACCENT3),
    ("Privacy Controls",
     "Users set profile to PUBLIC or PRIVATE. "
     "PRIVATE users excluded from search, explore, and suggestions. "
     "Privacy updated via PATCH /api/users/me/privacy. "
     "Three user types: STANDARD, BUSINESS, CREATOR.", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.3) + Inches(row * 1.9)
    add_card(slide7, x, y, Inches(5.9), Inches(1.7), title, desc, color)

# ══════════════════════════════════════════════════════════════
# SLIDE 8: DOCKER
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, C_WHITE)
add_header_bar(slide8, "Docker Containerization", "Database-per-Service pattern with Docker Compose")
add_slide_number(slide8, 8)

# Docker overview
add_text(slide8, Inches(0.5), Inches(1.3), Inches(12), Inches(0.3),
         "Docker Compose Architecture", size=18, color=C_DARK, bold=True)

# Docker containers visualization
containers = [
    ("mysql-user", "3307", "user_db", "Users, Auth\nSettings, Profiles", C_PRIMARY),
    ("mysql-post", "3308", "post_db", "Posts, Media\nScheduled, CTA", C_ACCENT1),
    ("mysql-interaction", "3309", "interaction_db", "Likes, Comments\nShares, Views", C_ACCENT2),
    ("mysql-connection", "3310", "connection_db", "Followers\nFollowing", C_SECONDARY),
    ("mysql-notification", "3311", "notification_db", "Notifications\nPreferences", C_ACCENT3),
]

for i, (name, port, db, desc, color) in enumerate(containers):
    x = Inches(0.3) + Inches(i * 2.6)
    y = Inches(1.8)
    # Container box
    add_rounded_rect(slide8, x, y, Inches(2.3), Inches(2.2), RGBColor(241, 245, 249))
    add_rect(slide8, x, y, Inches(2.3), Inches(0.08), color)
    # Docker icon area
    add_rounded_rect(slide8, x + Inches(0.6), y + Inches(0.2), Inches(1.1), Inches(0.5), color)
    add_text(slide8, x + Inches(0.6), y + Inches(0.22), Inches(1.1), Inches(0.5),
             "MySQL 8", size=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide8, x + Inches(0.15), y + Inches(0.8), Inches(2.0), Inches(0.3),
             name, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide8, x + Inches(0.15), y + Inches(1.1), Inches(2.0), Inches(0.25),
             f"Port: {port} | DB: {db}", size=9, color=C_MID, align=PP_ALIGN.CENTER)
    add_text(slide8, x + Inches(0.15), y + Inches(1.4), Inches(2.0), Inches(0.6),
             desc, size=9, color=C_MID, align=PP_ALIGN.CENTER)

# Docker benefits
add_text(slide8, Inches(0.5), Inches(4.3), Inches(12), Inches(0.3),
         "Why Docker for This Project?", size=18, color=C_DARK, bold=True)

benefits = [
    ("Isolation", "Each service has its own database\ncontainer — no schema conflicts.\nData persisted in Docker volumes.", C_PRIMARY),
    ("Reproducibility", "docker-compose up -d starts all\n5 databases with correct configs.\nSame environment everywhere.", C_ACCENT1),
    ("Database per Service", "True microservices pattern — each\nservice owns its data. No shared\ndatabase coupling between services.", C_ACCENT2),
    ("Easy Scaling", "Scale individual containers\nindependently based on load.\nAdd replicas as needed.", C_SECONDARY),
    ("Dev/Prod Parity", "Same Docker containers in dev,\nstaging, and production.\nNo 'works on my machine' issues.", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(benefits):
    x = Inches(0.3) + Inches(i * 2.6)
    y = Inches(4.7)
    add_card(slide8, x, y, Inches(2.3), Inches(1.8), title, desc, color)

# docker-compose snippet
add_rounded_rect(slide8, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), RGBColor(30, 41, 59))
add_text(slide8, Inches(0.7), Inches(6.73), Inches(12), Inches(0.4),
         "$ docker-compose up -d   |   5 MySQL containers   |   Persistent volumes   |   Auto-restart on failure   |   Network isolation",
         size=10, color=RGBColor(167, 243, 208))

# ══════════════════════════════════════════════════════════════
# SLIDE 9: CONCLUSION
# ══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, C_WHITE)
add_header_bar(slide9, "What Makes RevConnect Unique?", "How we stand apart from Instagram, Facebook & other social platforms")
add_slide_number(slide9, 9)

# Comparison table
add_text(slide9, Inches(0.5), Inches(1.3), Inches(12), Inches(0.3),
         "RevConnect vs Traditional Social Media Platforms", size=18, color=C_DARK, bold=True)

comparisons = [
    ("Microservices Architecture",
     "Instagram & Facebook use monolithic or large-scale proprietary architectures. RevConnect demonstrates a "
     "clean microservices approach with 7 independent services, each with its own database — making it "
     "highly scalable, maintainable, and fault-tolerant.", C_PRIMARY),
    ("Three User Types with Business Tools",
     "Unlike Instagram/Facebook which have limited business profiles, RevConnect offers Standard, Business, "
     "and Creator accounts — each with tailored features like CTA buttons (Shop Now, Learn More), paid "
     "partnership disclosures, product tagging, and post scheduling.", C_SECONDARY),
    ("Real Analytics with Multi-Service Aggregation",
     "RevConnect's analytics dashboard aggregates data from Interaction, Connection, and Post services "
     "via Feign clients — providing real engagement metrics, follower demographics by account type, "
     "and per-post performance. Not mock data — real aggregated microservice data.", C_ACCENT1),
    ("Granular Notification Preferences",
     "Users control notifications per type (likes, comments, shares, followers) independently — not "
     "just a global on/off switch. The notification service checks user preferences before creating "
     "any notification, respecting user choice at every level.", C_ACCENT2),
    ("True Service Discovery & Gateway",
     "Uses Netflix Eureka for dynamic service registration and Spring Cloud Gateway for intelligent "
     "routing, JWT authentication, and circuit breakers — a production-grade microservices setup "
     "that Instagram/Facebook don't expose to developers.", C_ACCENT4),
    ("Open Source & Educational",
     "Built as a learning project that demonstrates full-stack microservices development with "
     "real-world patterns: database-per-service, API gateway, Feign clients, Docker containers, "
     "CI/CD pipelines, and AWS deployment — a complete portfolio project.", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(comparisons):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.7) + Inches(row * 1.8)
    add_card(slide9, x, y, Inches(5.9), Inches(1.6), title, desc, color)

# ══════════════════════════════════════════════════════════════
# SLIDE 10: THANK YOU
# ══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, C_WHITE)

# Large decorative circles
add_circle(slide10, Inches(0.5), Inches(0.5), Inches(3), RGBColor(238, 242, 255))
add_circle(slide10, Inches(9.5), Inches(4.5), Inches(3.5), RGBColor(245, 243, 255))
add_circle(slide10, Inches(-1), Inches(5), Inches(2.5), RGBColor(236, 253, 245))
add_circle(slide10, Inches(8), Inches(-0.5), Inches(2), RGBColor(254, 243, 199))

# Center content
add_text(slide10, Inches(2), Inches(1.8), Inches(9.333), Inches(0.8),
         "Thank You!", size=60, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide10, Inches(2), Inches(3.0), Inches(9.333), Inches(0.5),
         "RevConnect — A Microservices-Based Social Media Platform", size=22, color=C_SECONDARY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide10, Inches(3), Inches(4.0), Inches(7.333), Inches(1.5),
         "Built with Angular 19 + Spring Boot 3 + Docker + AWS\n"
         "7 Microservices | 5 MySQL Databases | JWT Authentication\n"
         "Real-time Feed | Analytics | Messages | Business Tools",
         size=16, color=C_MID, align=PP_ALIGN.CENTER)

# Decorative line
add_rect(slide10, Inches(5), Inches(5.5), Inches(3.333), Inches(0.04), C_PRIMARY)

add_text(slide10, Inches(2), Inches(5.8), Inches(9.333), Inches(0.5),
         "Questions & Discussion", size=20, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide10, Inches(2), Inches(6.4), Inches(9.333), Inches(0.5),
         "We'd love to hear your feedback and answer any questions about the architecture and implementation.",
         size=14, color=C_MID, align=PP_ALIGN.CENTER)

add_slide_number(slide10, 10)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = r'd:\Microservices\RevConnect_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")

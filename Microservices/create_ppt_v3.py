from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ══════════════════════════════════════════════════════════════
# COLOR PALETTE
# ══════════════════════════════════════════════════════════════
C_PRIMARY = RGBColor(59, 130, 246)
C_SECONDARY = RGBColor(139, 92, 246)
C_ACCENT1 = RGBColor(16, 185, 129)
C_ACCENT2 = RGBColor(245, 158, 11)
C_ACCENT3 = RGBColor(239, 68, 68)
C_ACCENT4 = RGBColor(6, 182, 212)
C_DARK = RGBColor(30, 41, 59)
C_MID = RGBColor(71, 85, 105)
C_WHITE = RGBColor(255, 255, 255)
TOTAL_SLIDES = 14

def add_bg(slide, color=C_WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=C_DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body, icon_color, title_color=C_DARK):
    add_rounded_rect(slide, left, top, width, height, C_WHITE)
    add_rect(slide, left, top, width, Inches(0.06), icon_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.35),
             title, size=13, color=title_color, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
             body, size=10, color=C_MID)

def add_slide_number(slide, num):
    add_text(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
             f"{num}/{TOTAL_SLIDES}", size=10, color=C_MID, align=PP_ALIGN.RIGHT)

def add_header_bar(slide, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), C_PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
             title, size=28, color=C_WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4),
                 subtitle, size=14, color=RGBColor(191, 219, 254))
    add_circle(slide, Inches(11.5), Inches(-0.3), Inches(0.8), RGBColor(96, 165, 250))
    add_circle(slide, Inches(12.2), Inches(0.3), Inches(0.6), C_SECONDARY)


# ══════════════════════════════════════════════════════════════
# SLIDE 1: DEPLOYMENT (AWS + Jenkins CI/CD)
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, C_WHITE)
add_header_bar(slide1, "Deployment Pipeline", "AWS Cloud Deployment with Jenkins CI/CD for Monolithic Version")
add_slide_number(slide1, 1)

stages = [
    ("Developer\nPushes Code", "Git commit & push\nto GitHub repo", RGBColor(99, 102, 241)),
    ("GitHub\nWebhook", "Triggers Jenkins\nbuild pipeline", C_DARK),
    ("Jenkins\nCI Pipeline", "Maven build, run\nunit tests, package", C_ACCENT2),
    ("Docker\nBuild", "Dockerfile builds\napp container image", C_ACCENT4),
    ("Push to\nDocker Hub / ECR", "Image pushed to\ncontainer registry", C_SECONDARY),
    ("AWS EC2\nDeploy", "Jenkins deploys to\nEC2 via SSH/ECS", C_ACCENT1),
]
for i, (title, desc, color) in enumerate(stages):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(1.4)
    add_rounded_rect(slide1, x, y, Inches(1.95), Inches(1.6), color)
    add_text(slide1, x + Inches(0.1), y + Inches(0.15), Inches(1.75), Inches(0.5),
             title, size=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, x + Inches(0.1), y + Inches(0.7), Inches(1.75), Inches(0.7),
             desc, size=10, color=RGBColor(220, 230, 255), align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_arrow(slide1, x + Inches(1.95), y + Inches(0.6), Inches(0.2), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide1, Inches(0.5), Inches(3.3), Inches(12), Inches(0.4),
         "AWS Cloud Architecture", size=18, color=C_DARK, bold=True)

aws_components = [
    ("Route 53 / DNS", "Domain routing to\nload balancer", RGBColor(99, 102, 241)),
    ("Elastic Load\nBalancer (ALB)", "Distributes traffic\nacross instances", C_ACCENT2),
    ("EC2 Instance", "Runs Spring Boot\napp in Docker", C_ACCENT1),
    ("RDS MySQL", "Managed MySQL\ndatabase service", C_PRIMARY),
    ("S3 Bucket", "Static assets &\nmedia file storage", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(aws_components):
    x = Inches(0.3) + Inches(i * 2.6)
    y = Inches(3.8)
    add_rounded_rect(slide1, x, y, Inches(2.3), Inches(1.2), C_WHITE)
    add_rect(slide1, x, y, Inches(2.3), Inches(0.06), color)
    add_text(slide1, x + Inches(0.15), y + Inches(0.15), Inches(2.0), Inches(0.4),
             title, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, x + Inches(0.15), y + Inches(0.55), Inches(2.0), Inches(0.5),
             desc, size=9, color=C_MID, align=PP_ALIGN.CENTER)
    if i < len(aws_components) - 1:
        add_arrow(slide1, x + Inches(2.3), y + Inches(0.45), Inches(0.3), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide1, Inches(0.5), Inches(5.3), Inches(12), Inches(0.3),
         "Jenkins Pipeline Stages (Jenkinsfile)", size=14, color=C_DARK, bold=True)
jenkins_cards = [
    ("Checkout", "git checkout from\nGitHub repository", RGBColor(99, 102, 241)),
    ("Build", "mvn clean package\n-DskipTests", C_ACCENT2),
    ("Test", "mvn test \u2014 runs\nunit & integration tests", C_ACCENT1),
    ("Docker Build", "docker build -t\nrevconnect:latest .", C_ACCENT4),
    ("Push Image", "docker push to\nDocker Hub / AWS ECR", C_SECONDARY),
    ("Deploy", "ssh to EC2, pull image\ndocker-compose up -d", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(jenkins_cards):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(5.7)
    add_rounded_rect(slide1, x, y, Inches(1.95), Inches(1.1), C_WHITE)
    add_rect(slide1, x, y, Inches(1.95), Inches(0.05), color)
    add_circle(slide1, x + Inches(0.1), y + Inches(0.15), Inches(0.3), color)
    add_text(slide1, x + Inches(0.1), y + Inches(0.17), Inches(0.3), Inches(0.3),
             str(i + 1), size=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, x + Inches(0.5), y + Inches(0.12), Inches(1.3), Inches(0.25),
             title, size=10, color=color, bold=True)
    add_text(slide1, x + Inches(0.5), y + Inches(0.4), Inches(1.3), Inches(0.6),
             desc, size=9, color=C_MID)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: INTRODUCTION (Abstract + Tech Stack)
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, C_WHITE)

add_rect(slide2, Inches(0), Inches(0), Inches(5.5), prs.slide_height, RGBColor(238, 242, 255))
add_rect(slide2, Inches(0), Inches(0), Inches(0.08), prs.slide_height, C_PRIMARY)
add_circle(slide2, Inches(10), Inches(0.5), Inches(2.5), RGBColor(238, 242, 255))
add_circle(slide2, Inches(11), Inches(5), Inches(2), RGBColor(245, 243, 255))
add_circle(slide2, Inches(4.5), Inches(5.5), Inches(1.5), RGBColor(236, 253, 245))

add_text(slide2, Inches(0.8), Inches(0.8), Inches(4), Inches(0.5),
         "RevConnect", size=42, color=C_PRIMARY, bold=True)
add_text(slide2, Inches(0.8), Inches(1.4), Inches(4), Inches(0.5),
         "Social Media Platform", size=22, color=C_SECONDARY, bold=True)
add_text(slide2, Inches(0.8), Inches(2.2), Inches(4.2), Inches(0.3),
         "ABSTRACT", size=14, color=C_PRIMARY, bold=True)

abstract_text = (
    "RevConnect is a full-stack professional social media platform built using a microservices architecture "
    "with Angular 19 frontend and Spring Boot 3 backend. The system comprises 7 independent microservices "
    "(User, Post, Interaction, Connection, Notification, Feed, API Gateway) each with its own MySQL database, "
    "communicating via REST APIs and Eureka service discovery.\n\n"
    "Key features include user registration with email verification, JWT-based authentication, real-time "
    "personalized feeds, post creation with media uploads (images/videos), social interactions (likes, comments, "
    "shares), connection management (follow/unfollow), direct messaging with WebSocket support, push/email "
    "notifications with granular preferences, user analytics dashboards, stories, bookmarks, explore/search, "
    "and business tools (CTA buttons, paid partnerships, product tagging, scheduled posts).\n\n"
    "The platform supports three user types: Standard, Business, and Creator \u2014 each with tailored features. "
    "Privacy controls allow users to set profiles as Public or Private. The architecture ensures scalability, "
    "fault tolerance via circuit breakers, and clean separation of concerns."
)
add_text(slide2, Inches(0.8), Inches(2.6), Inches(4.2), Inches(4.5),
         abstract_text, size=11, color=C_MID)

add_text(slide2, Inches(6), Inches(0.8), Inches(6), Inches(0.4),
         "Technology Stack", size=20, color=C_DARK, bold=True)
techs = [
    ("Frontend", "Angular 19, TypeScript, RxJS, TailwindCSS, Font Awesome", C_PRIMARY),
    ("Backend", "Spring Boot 3, Java 17, Spring Cloud, Eureka, Feign, JWT", C_SECONDARY),
    ("Database", "MySQL 8 (5 Docker containers), JPA/Hibernate, Spring Data", C_ACCENT1),
    ("DevOps", "Docker Compose, API Gateway, Circuit Breakers, Eureka Discovery", C_ACCENT2),
    ("Features", "Real-time Feed, Analytics, Stories, Messages, Notifications, Business Tools", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(techs):
    y = Inches(1.4) + Inches(i * 1.1)
    add_card(slide2, Inches(6), y, Inches(6.8), Inches(0.9), title, desc, color)
add_slide_number(slide2, 2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, C_WHITE)
add_header_bar(slide3, "System Architecture", "Microservices Architecture with Service Discovery & API Gateway")
add_slide_number(slide3, 3)

add_rounded_rect(slide3, Inches(5.2), Inches(1.3), Inches(3), Inches(0.7), RGBColor(219, 234, 254))
add_text(slide3, Inches(5.2), Inches(1.35), Inches(3), Inches(0.35),
         "Angular 19 Frontend", size=14, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, Inches(5.2), Inches(1.65), Inches(3), Inches(0.3),
         "localhost:4200", size=10, color=C_MID, align=PP_ALIGN.CENTER)
add_arrow(slide3, Inches(6.4), Inches(2.1), Inches(0.6), Inches(0.3), C_MID)

add_rounded_rect(slide3, Inches(4.5), Inches(2.5), Inches(4.4), Inches(0.8), RGBColor(254, 243, 199))
add_text(slide3, Inches(4.5), Inches(2.55), Inches(4.4), Inches(0.4),
         "API Gateway (Spring Cloud Gateway)", size=14, color=C_ACCENT2, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, Inches(4.5), Inches(2.9), Inches(4.4), Inches(0.3),
         "JWT Auth | Route | Load Balance | Circuit Breaker \u2014 :8080", size=10, color=C_MID, align=PP_ALIGN.CENTER)

add_rounded_rect(slide3, Inches(10), Inches(2.5), Inches(2.8), Inches(0.8), RGBColor(237, 233, 254))
add_text(slide3, Inches(10), Inches(2.55), Inches(2.8), Inches(0.4),
         "Eureka Server", size=14, color=C_SECONDARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide3, Inches(10), Inches(2.9), Inches(2.8), Inches(0.3),
         "Service Discovery \u2014 :8761", size=10, color=C_MID, align=PP_ALIGN.CENTER)

services = [
    ("User\nService", ":8081", C_PRIMARY, "Auth, Profile\nSettings, Search"),
    ("Post\nService", ":8082", C_ACCENT1, "CRUD, Schedule\nMedia, CTA"),
    ("Feed\nService", ":8083", C_ACCENT4, "Personalized\nFeed Aggregation"),
    ("Interaction\nService", ":8084", C_ACCENT2, "Likes, Comments\nShares, Views"),
    ("Connection\nService", ":8085", C_SECONDARY, "Follow/Unfollow\nStats, Network"),
    ("Notification\nService", ":8086", C_ACCENT3, "Push, Email\nPreferences"),
]
for i, (name, port, color, desc) in enumerate(services):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(4.0)
    add_rounded_rect(slide3, x, y, Inches(1.95), Inches(1.8), C_WHITE)
    add_rect(slide3, x, y, Inches(1.95), Inches(0.08), color)
    add_text(slide3, x, y + Inches(0.15), Inches(1.95), Inches(0.5),
             name, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, x, y + Inches(0.65), Inches(1.95), Inches(0.25),
             port, size=9, color=C_MID, align=PP_ALIGN.CENTER)
    add_text(slide3, x, y + Inches(0.9), Inches(1.95), Inches(0.7),
             desc, size=9, color=C_MID, align=PP_ALIGN.CENTER)

dbs = [
    ("user_db", ":3307", C_PRIMARY),
    ("post_db", ":3308", C_ACCENT1),
    ("\u2014", "", C_ACCENT4),
    ("interaction_db", ":3309", C_ACCENT2),
    ("connection_db", ":3310", C_SECONDARY),
    ("notification_db", ":3311", C_ACCENT3),
]
for i, (name, port, color) in enumerate(dbs):
    if name == "\u2014": continue
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(6.1)
    add_rounded_rect(slide3, x, y, Inches(1.95), Inches(0.7), RGBColor(241, 245, 249))
    add_text(slide3, x, y + Inches(0.05), Inches(1.95), Inches(0.35),
             f"MySQL {name}", size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, x, y + Inches(0.35), Inches(1.95), Inches(0.25),
             f"Docker {port}", size=8, color=C_MID, align=PP_ALIGN.CENTER)

add_text(slide3, Inches(0.3), Inches(6.9), Inches(12), Inches(0.3),
         "Each microservice has its own isolated MySQL database running in Docker containers \u2014 Database per Service pattern",
         size=10, color=C_MID, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: LOGIN, REGISTER, FORGOT PASSWORD
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, C_WHITE)
add_header_bar(slide4, "Authentication Flow", "Login, Register & Forgot Password")
add_slide_number(slide4, 4)

add_text(slide4, Inches(0.5), Inches(1.3), Inches(5), Inches(0.4),
         "Login Flow", size=18, color=C_PRIMARY, bold=True)
login_steps = [
    ("1", "User enters\nemail & password", C_PRIMARY),
    ("2", "Angular sends\nPOST /api/auth/login", C_ACCENT4),
    ("3", "Gateway routes\nto User Service", C_ACCENT2),
    ("4", "BCrypt verifies\npassword hash", C_ACCENT3),
    ("5", "JWT token\ngenerated & returned", C_ACCENT1),
]
for i, (num, text, color) in enumerate(login_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_circle(slide4, x, Inches(1.9), Inches(0.5), color)
    add_text(slide4, x, Inches(1.93), Inches(0.5), Inches(0.5),
             num, size=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide4, x - Inches(0.3), Inches(2.5), Inches(1.1), Inches(0.6),
             text, size=10, color=C_DARK, align=PP_ALIGN.CENTER)
    if i < len(login_steps) - 1:
        add_arrow(slide4, x + Inches(0.6), Inches(2.05), Inches(1.6), Inches(0.2), RGBColor(203, 213, 225))

add_text(slide4, Inches(0.5), Inches(3.2), Inches(12), Inches(0.5),
         "Dependencies: BCryptPasswordEncoder, JwtTokenProvider (HS256, 256-bit secret), Spring Security FilterChain",
         size=10, color=C_MID)

add_text(slide4, Inches(0.5), Inches(3.8), Inches(5), Inches(0.4),
         "Register Flow", size=18, color=C_ACCENT1, bold=True)
reg_steps = [
    ("1", "User fills form:\nname, email, password", C_ACCENT1),
    ("2", "POST /api/auth/\nregister", C_ACCENT4),
    ("3", "Duplicate check\n+ BCrypt hash", C_ACCENT2),
    ("4", "Verification email\nsent (code in DB)", C_ACCENT3),
    ("5", "User clicks link\naccount activated", C_PRIMARY),
]
for i, (num, text, color) in enumerate(reg_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_circle(slide4, x, Inches(4.4), Inches(0.5), color)
    add_text(slide4, x, Inches(4.43), Inches(0.5), Inches(0.5),
             num, size=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide4, x - Inches(0.3), Inches(5.0), Inches(1.1), Inches(0.6),
             text, size=10, color=C_DARK, align=PP_ALIGN.CENTER)
    if i < len(reg_steps) - 1:
        add_arrow(slide4, x + Inches(0.6), Inches(4.55), Inches(1.6), Inches(0.2), RGBColor(203, 213, 225))

add_text(slide4, Inches(0.5), Inches(5.8), Inches(5), Inches(0.4),
         "Forgot Password Flow", size=18, color=C_ACCENT3, bold=True)
add_text(slide4, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
         "User enters email \u2192 POST /api/auth/forgot-password \u2192 User Service generates reset token (UUID) \u2192 "
         "Email sent with reset link \u2192 User clicks link with token \u2192 Enters new password \u2192 "
         "POST /api/auth/reset-password with token \u2192 Password re-hashed with BCrypt \u2192 Success",
         size=11, color=C_MID)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: HOME PAGE & FEED
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, C_WHITE)
add_header_bar(slide5, "Home Page & Personalized Feed", "How posts are loaded and displayed")
add_slide_number(slide5, 5)

flow_items = [
    ("User opens /feed", "FeedPage component\ninitializes", C_PRIMARY),
    ("Load profile", "GET /api/users/me\nvia API Gateway + JWT", C_ACCENT4),
    ("Fetch feed", "GET /api/posts/feed/\npersonalized", C_ACCENT1),
    ("Feed Service", "Gets following IDs from\nConnection Service (Feign)", C_ACCENT2),
    ("Post Service", "Queries posts by\nfollowing user IDs", C_SECONDARY),
    ("Enrich posts", "Fetches author info\nfrom User Service", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(flow_items):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(1.4)
    add_rounded_rect(slide5, x, y, Inches(1.95), Inches(1.5), C_WHITE)
    add_rect(slide5, x, y, Inches(1.95), Inches(0.06), color)
    add_text(slide5, x + Inches(0.1), y + Inches(0.15), Inches(1.75), Inches(0.3),
             title, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide5, x + Inches(0.1), y + Inches(0.5), Inches(1.75), Inches(0.8),
             desc, size=10, color=C_MID, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_arrow(slide5, x + Inches(1.95), y + Inches(0.6), Inches(0.2), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide5, Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
         "Feed Page Components & Features", size=16, color=C_DARK, bold=True)
cards_data = [
    ("Stories Feed", "Horizontal scrollable story cards at top.\nUsers can create and view 24h stories.", C_PRIMARY),
    ("Create Post Box", "Text input + Media upload + Schedule +\nCTA buttons. Supports IMAGE/VIDEO/TEXT types.", C_ACCENT1),
    ("Live Analytics Sidebar", "Shows total views, likes, shares, comments\nin real-time for Business/Creator users.", C_ACCENT2),
    ("Trending Topics", "Fetched from search service. Shows\npopular hashtags and search terms.", C_SECONDARY),
    ("Post Cards", "Each post shows author info, content, media,\nlike/comment/share buttons, timestamps.", C_ACCENT3),
    ("Infinite Scroll", "Paginated feed (page=0, size=10).\nLoads more posts as user scrolls.", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(cards_data):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(3.7) + Inches(row * 1.7)
    add_card(slide5, x, y, Inches(3.8), Inches(1.5), title, desc, color)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: POSTS \u2014 CREATE, SCHEDULE, CTA
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, C_WHITE)
add_header_bar(slide6, "Post Creation & Business Tools", "Create, Schedule, CTA, Partnerships & Product Tags")
add_slide_number(slide6, 6)

add_text(slide6, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
         "Normal Post Creation", size=16, color=C_PRIMARY, bold=True)
post_steps = [
    ("Write content\n+ attach media", C_PRIMARY),
    ("Upload media to\nPOST /api/media/upload", C_ACCENT4),
    ("POST /api/posts\nwith content + mediaUrls", C_ACCENT1),
    ("PostService saves\nto MySQL (post_db)", C_ACCENT2),
    ("Post appears\nin feed", C_SECONDARY),
]
for i, (text, color) in enumerate(post_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_rounded_rect(slide6, x, Inches(1.7), Inches(2.1), Inches(0.8), color)
    add_text(slide6, x + Inches(0.1), Inches(1.75), Inches(1.9), Inches(0.7),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(post_steps) - 1:
        add_arrow(slide6, x + Inches(2.1), Inches(1.95), Inches(0.4), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide6, Inches(0.5), Inches(2.8), Inches(5), Inches(0.3),
         "Scheduled Post with CTA", size=16, color=C_ACCENT2, bold=True)
sched_steps = [
    ("User selects date,\ntime, CTA label & URL", C_ACCENT2),
    ("Frontend builds ISO\ndatetime + request body", C_ACCENT4),
    ("POST /api/posts/\nschedule", C_ACCENT1),
    ("Post saved with\nisPublished=false", C_ACCENT3),
    ("Scheduler publishes\nat scheduled time", C_SECONDARY),
]
for i, (text, color) in enumerate(sched_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_rounded_rect(slide6, x, Inches(3.2), Inches(2.1), Inches(0.8), color)
    add_text(slide6, x + Inches(0.1), Inches(3.25), Inches(1.9), Inches(0.7),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(sched_steps) - 1:
        add_arrow(slide6, x + Inches(2.1), Inches(3.45), Inches(0.4), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide6, Inches(0.5), Inches(4.3), Inches(5), Inches(0.3),
         "User Type Differences", size=16, color=C_DARK, bold=True)
add_card(slide6, Inches(0.5), Inches(4.7), Inches(5.8), Inches(2.4),
         "BUSINESS Account",
         "\u2022 Post Categories: Standard, Announcement, Product Update\n"
         "\u2022 Call To Action: Shop Now, Learn More, Sign Up, Contact Us,\n  Download, Subscribe, Book Now\n"
         "\u2022 CTA URL: Link to external page\n"
         "\u2022 No paid partnership or product tagging",
         C_PRIMARY)
add_card(slide6, Inches(6.8), Inches(4.7), Inches(6), Inches(2.4),
         "CREATOR Account",
         "\u2022 All Business features PLUS:\n"
         "\u2022 Paid Partnership Disclosure: Mark as sponsored,\n  enter brand partner name (e.g. Nike, Apple)\n"
         "\u2022 Product Tagging: Comma-separated tags\n  (e.g. Shoes, Watch, Sunglasses)\n"
         "\u2022 Partnership info stored in partnerName & productTags fields",
         C_SECONDARY)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: LIKES, COMMENTS, SHARES
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, C_WHITE)
add_header_bar(slide7, "Social Interactions", "Likes, Comments & Shares \u2014 Interaction Service")
add_slide_number(slide7, 7)

add_text(slide7, Inches(0.5), Inches(1.3), Inches(4), Inches(0.3),
         "Like Flow", size=16, color=C_ACCENT3, bold=True)
add_text(slide7, Inches(0.5), Inches(1.7), Inches(12), Inches(1.0),
         "User clicks \u2764 \u2192 Frontend calls POST /api/posts/{id}/like \u2192 "
         "API Gateway routes to Interaction Service \u2192 "
         "Checks if already liked (unique constraint: user_id + post_id) \u2192 "
         "If not liked: saves Like entity, sends notification to post author via Notification Service (Feign) \u2192 "
         "If already liked: returns 'Already liked'. "
         "Unlike: DELETE /api/posts/{id}/like removes the Like record.",
         size=11, color=C_MID)

add_text(slide7, Inches(0.5), Inches(2.8), Inches(4), Inches(0.3),
         "Comment Flow", size=16, color=C_PRIMARY, bold=True)
add_text(slide7, Inches(0.5), Inches(3.2), Inches(12), Inches(1.0),
         "User types comment \u2192 POST /api/posts/{id}/comments with {content, parentId?} \u2192 "
         "Interaction Service saves Comment entity (supports nested replies via parentId) \u2192 "
         "Notification sent to post author \u2192 Comments fetched with pagination: "
         "GET /api/posts/{id}/comments?page=0&size=10 \u2192 "
         "Each comment enriched with author name/avatar from User Service.",
         size=11, color=C_MID)

add_text(slide7, Inches(0.5), Inches(4.3), Inches(4), Inches(0.3),
         "Share Flow", size=16, color=C_ACCENT1, bold=True)
add_text(slide7, Inches(0.5), Inches(4.7), Inches(12), Inches(1.0),
         "User clicks Share \u2192 POST /api/posts/{id}/share with optional {comment} \u2192 "
         "Creates Share entity with REPOST type \u2192 Notification sent to original post author \u2192 "
         "Share count visible in post cards and analytics. "
         "Shares tracked per user in shares table.",
         size=11, color=C_MID)

add_text(slide7, Inches(0.5), Inches(5.6), Inches(4), Inches(0.3),
         "View Tracking", size=16, color=C_ACCENT4, bold=True)
add_text(slide7, Inches(0.5), Inches(6.0), Inches(12), Inches(1.0),
         "When a user's feed loads, each post records a view: POST /api/posts/{id}/view \u2192 "
         "PostView entity with unique constraint (post_id + user_id) ensures each user counted only once per post \u2192 "
         "View count used in analytics dashboard. No duplicate counting on page refresh.",
         size=11, color=C_MID)

add_rounded_rect(slide7, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), RGBColor(241, 245, 249))
add_text(slide7, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "Dependencies: LikeRepository, CommentRepository, ShareRepository, PostViewRepository, NotificationServiceClient (Feign), UserServiceClient (Feign)",
         size=10, color=C_MID)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: MESSAGES
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, C_WHITE)
add_header_bar(slide8, "Direct Messages", "Real-time messaging between connected users")
add_slide_number(slide8, 8)

msg_cards = [
    ("Conversation List", "GET /api/messages/conversations\nReturns list of conversations with last message preview,\ntimestamp, unread count. Sorted by newest first.", C_PRIMARY),
    ("Load Messages", "GET /api/messages/conversation/{recipientId}\nFetches paginated messages between two users.\nMessages shown in chronological order.", C_ACCENT4),
    ("Send Message", "POST /api/messages with {recipientId, content, type}\nSupports TEXT, IMAGE, VOICE message types.\nNotification sent to recipient.", C_ACCENT1),
    ("Real-time Updates", "WebSocket connection via STOMP protocol.\nNew messages pushed instantly to recipients.\nNo polling needed \u2014 server pushes updates.", C_ACCENT2),
    ("Message Features", "Read receipts, typing indicators, message deletion.\nVoice messages recorded in browser and uploaded.\nImage messages use media upload service.", C_SECONDARY),
    ("Architecture", "Messages stored in User Service database.\nFrontend uses RxJS observables for reactive updates.\nConversations sorted by lastMessageTime descending.", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(msg_cards):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.4) + Inches(row * 2.8)
    add_card(slide8, x, y, Inches(3.8), Inches(2.5), title, desc, color)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: EXPLORE & CONNECTIONS
# ══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, C_WHITE)
add_header_bar(slide9, "Explore & Connections", "Search, Follow/Unfollow & Network Management")
add_slide_number(slide9, 9)

add_text(slide9, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
         "Explore / Search", size=16, color=C_PRIMARY, bold=True)
search_steps = [
    ("User types query\nin search bar", C_PRIMARY),
    ("GET /api/search/all?\nquery=... (debounced)", C_ACCENT4),
    ("SearchController calls\nUserService + PostService", C_ACCENT1),
    ("Users filtered: PRIVATE\naccounts excluded", C_ACCENT3),
    ("Results shown in\nAll / Users / Posts tabs", C_SECONDARY),
]
for i, (text, color) in enumerate(search_steps):
    x = Inches(0.3) + Inches(i * 2.55)
    add_rounded_rect(slide9, x, Inches(1.7), Inches(2.2), Inches(0.9), color)
    add_text(slide9, x + Inches(0.1), Inches(1.75), Inches(2.0), Inches(0.8),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(search_steps) - 1:
        add_arrow(slide9, x + Inches(2.2), Inches(2.0), Inches(0.35), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide9, Inches(0.5), Inches(2.9), Inches(5), Inches(0.3),
         "Follow / Unfollow Flow", size=16, color=C_ACCENT1, bold=True)
conn_steps = [
    ("User clicks Follow\non profile card", C_ACCENT1),
    ("POST /api/users/\n{id}/follow", C_ACCENT4),
    ("Connection Service\ncreates Follow record", C_ACCENT2),
    ("FOLLOW notification\nsent to target user", C_ACCENT3),
    ("Follower/Following\ncounts updated", C_PRIMARY),
]
for i, (text, color) in enumerate(conn_steps):
    x = Inches(0.3) + Inches(i * 2.55)
    add_rounded_rect(slide9, x, Inches(3.3), Inches(2.2), Inches(0.9), color)
    add_text(slide9, x + Inches(0.1), Inches(3.35), Inches(2.0), Inches(0.8),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(conn_steps) - 1:
        add_arrow(slide9, x + Inches(2.2), Inches(3.6), Inches(0.35), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide9, Inches(0.5), Inches(4.5), Inches(5), Inches(0.3),
         "Privacy & Suggestions", size=16, color=C_DARK, bold=True)
priv_cards = [
    ("Privacy Controls",
     "Users can set profile to PUBLIC or PRIVATE.\nPRIVATE users are excluded from:\n\u2022 Search results (searchUsers query filter)\n\u2022 Suggested users (findSuggestions query filter)\n\u2022 Explore page discovery\nPrivacy updated via PATCH /api/users/me/privacy",
     C_ACCENT3),
    ("Suggested Users",
     "GET /api/users/suggested returns non-private,\nemail-verified users excluding current user.\nShown in Explore page sidebar.\nAlgorithm: Basic exclusion filter.\nFuture: ML-based recommendation engine.",
     C_PRIMARY),
    ("Connection Stats",
     "GET /api/users/{id}/connection-stats returns:\n\u2022 followersCount \u2014 users following you\n\u2022 followingCount \u2014 users you follow\nUsed in profile page and analytics.\nStored in Connection Service (connection_db).",
     C_ACCENT1),
]
for i, (title, desc, color) in enumerate(priv_cards):
    x = Inches(0.5) + Inches(i * 4.2)
    add_card(slide9, x, Inches(4.9), Inches(3.8), Inches(2.4), title, desc, color)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: PROFILE & NOTIFICATIONS
# ══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, C_WHITE)
add_header_bar(slide10, "Profile & Notifications", "User profile management and notification system")
add_slide_number(slide10, 10)

add_text(slide10, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
         "Profile Page Features", size=16, color=C_PRIMARY, bold=True)
profile_cards = [
    ("Profile Display", "Cover photo, avatar, name, bio, location,\nwebsite, user type badge, follower stats.\nGET /api/users/{id} + connection-stats", C_PRIMARY),
    ("Edit Profile", "PUT /api/users/me with ProfileUpdateRequest.\nUpdates name, bio, location, avatar,\ncover photo, business info, social links.", C_ACCENT1),
    ("User Posts Tab", "GET /api/posts/user/{id} \u2014 paginated.\nShows all posts by the user with\nlike/comment/share counts.", C_ACCENT2),
    ("Bookmarks", "GET /api/interactions/bookmarks.\nSaved posts for later reading.\nToggle via POST/DELETE bookmark endpoint.", C_SECONDARY),
]
for i, (title, desc, color) in enumerate(profile_cards):
    x = Inches(0.5) + Inches(i * 3.15)
    add_card(slide10, x, Inches(1.7), Inches(2.9), Inches(1.8), title, desc, color)

add_text(slide10, Inches(0.5), Inches(3.8), Inches(5), Inches(0.3),
         "Notification System", size=16, color=C_ACCENT3, bold=True)
notif_steps = [
    ("Action triggers\nnotification", C_PRIMARY),
    ("Interaction Service\ncalls Notification API", C_ACCENT4),
    ("Check user prefs\n(granular per-type)", C_ACCENT2),
    ("If allowed: save\nNotification entity", C_ACCENT1),
    ("Frontend polls\nunread count", C_ACCENT3),
]
for i, (text, color) in enumerate(notif_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_rounded_rect(slide10, x, Inches(4.2), Inches(2.1), Inches(0.8), color)
    add_text(slide10, x + Inches(0.1), Inches(4.25), Inches(1.9), Inches(0.7),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(notif_steps) - 1:
        add_arrow(slide10, x + Inches(2.1), Inches(4.45), Inches(0.4), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide10, Inches(0.5), Inches(5.3), Inches(12), Inches(0.3),
         "Granular Notification Preferences", size=14, color=C_DARK, bold=True)
pref_items = [
    ("Likes", "notifyLike", C_ACCENT3),
    ("Comments", "notifyComment", C_PRIMARY),
    ("New Followers", "notifyNewFollower", C_ACCENT1),
    ("Shares", "notifyShare", C_ACCENT4),
    ("Connections", "notifyConnectionRequest", C_ACCENT2),
    ("Email", "emailNotifications", C_SECONDARY),
]
for i, (title, pref, color) in enumerate(pref_items):
    x = Inches(0.5) + Inches(i * 2.1)
    add_rounded_rect(slide10, x, Inches(5.7), Inches(1.9), Inches(0.8), C_WHITE)
    add_rect(slide10, x, Inches(5.7), Inches(1.9), Inches(0.05), color)
    add_text(slide10, x + Inches(0.1), Inches(5.8), Inches(1.7), Inches(0.3),
             title, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide10, x + Inches(0.1), Inches(6.1), Inches(1.7), Inches(0.3),
             pref, size=9, color=C_MID, align=PP_ALIGN.CENTER)

add_text(slide10, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
         "Types: LIKE, COMMENT, SHARE, FOLLOW, MENTION, SYSTEM | Checked in NotificationController before saving | Stored in UserSettings entity",
         size=10, color=C_MID)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: ANALYTICS
# ══════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11, C_WHITE)
add_header_bar(slide11, "Analytics Dashboard", "Data aggregation from multiple services")
add_slide_number(slide11, 11)

add_text(slide11, Inches(0.5), Inches(1.3), Inches(12), Inches(0.3),
         "How Analytics Collects Data", size=16, color=C_PRIMARY, bold=True)
analytics_flow = [
    ("User opens\nAnalytics page", C_PRIMARY),
    ("GET /api/analytics/\noverview", C_ACCENT4),
    ("Fetch follower count\nfrom Connection Service", C_ACCENT1),
    ("Fetch user's post IDs\nfrom Post Service", C_ACCENT2),
    ("Count likes, comments,\nshares, views per post", C_ACCENT3),
    ("Aggregate &\nreturn totals", C_SECONDARY),
]
for i, (text, color) in enumerate(analytics_flow):
    x = Inches(0.3) + Inches(i * 2.15)
    add_rounded_rect(slide11, x, Inches(1.7), Inches(1.95), Inches(0.9), color)
    add_text(slide11, x + Inches(0.1), Inches(1.75), Inches(1.75), Inches(0.8),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(analytics_flow) - 1:
        add_arrow(slide11, x + Inches(1.95), Inches(2.0), Inches(0.2), Inches(0.15), RGBColor(203, 213, 225))

add_text(slide11, Inches(0.5), Inches(2.9), Inches(12), Inches(0.3),
         "Analytics Endpoints & Metrics", size=16, color=C_DARK, bold=True)
metrics = [
    ("Overview", "/api/analytics/overview\n\nReturns: totalViews, totalLikes,\ntotalComments, totalShares,\ntotalFollowers, totalPosts", C_PRIMARY),
    ("Post Performance", "/api/analytics/posts\n\nPer-post breakdown: likes, comments,\nshares, views for each post.\nTop 5 posts by engagement.", C_ACCENT1),
    ("Follower Growth", "/api/analytics/followers\n\nFollower count over time.\nNew followers per day/week.\nGrowth trend data.", C_ACCENT2),
    ("Engagement Rate", "/api/analytics/engagement\n\nTotal engagement across all posts.\nEngagement = likes + comments + shares.\nRate = engagement / views.", C_ACCENT3),
    ("Audience Demographics", "/api/analytics/audience\n\nReal follower breakdown by account type.\nPersonal, Creator, Business counts.\nFetched from User Service via Feign.", C_SECONDARY),
    ("Profile Views", "/api/analytics/profile-views\n\nTrack who views your profile.\nReal view counting with\nPostView entity (unique per user).", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(3.3) + Inches(row * 2.1)
    add_card(slide11, x, y, Inches(3.8), Inches(1.9), title, desc, color)

add_rounded_rect(slide11, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35), RGBColor(241, 245, 249))
add_text(slide11, Inches(0.7), Inches(7.03), Inches(12), Inches(0.3),
         "Dependencies: AnalyticsController (Interaction Service) \u2192 PostServiceClient (Feign) \u2192 ConnectionServiceClient (Feign) \u2192 LikeRepo, CommentRepo, ShareRepo, PostViewRepo",
         size=9, color=C_MID)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: DOCKER
# ══════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12, C_WHITE)
add_header_bar(slide12, "Docker Containerization", "Database-per-Service pattern with Docker Compose")
add_slide_number(slide12, 12)

add_text(slide12, Inches(0.5), Inches(1.3), Inches(12), Inches(0.3),
         "Docker Compose Architecture", size=18, color=C_DARK, bold=True)

containers = [
    ("mysql-user", "3307", "user_db", "Users, Auth\nSettings, Profiles", C_PRIMARY),
    ("mysql-post", "3308", "post_db", "Posts, Media\nScheduled, CTA", C_ACCENT1),
    ("mysql-interaction", "3309", "interaction_db", "Likes, Comments\nShares, Views", C_ACCENT2),
    ("mysql-connection", "3310", "connection_db", "Followers\nFollowing", C_SECONDARY),
    ("mysql-notification", "3311", "notification_db", "Notifications\nPreferences", C_ACCENT3),
]
for i, (name, port, db, desc, color) in enumerate(containers):
    x = Inches(0.3) + Inches(i * 2.6)
    y = Inches(1.8)
    add_rounded_rect(slide12, x, y, Inches(2.3), Inches(2.2), RGBColor(241, 245, 249))
    add_rect(slide12, x, y, Inches(2.3), Inches(0.08), color)
    add_rounded_rect(slide12, x + Inches(0.6), y + Inches(0.2), Inches(1.1), Inches(0.5), color)
    add_text(slide12, x + Inches(0.6), y + Inches(0.22), Inches(1.1), Inches(0.5),
             "MySQL 8", size=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide12, x + Inches(0.15), y + Inches(0.8), Inches(2.0), Inches(0.3),
             name, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide12, x + Inches(0.15), y + Inches(1.1), Inches(2.0), Inches(0.25),
             f"Port: {port} | DB: {db}", size=9, color=C_MID, align=PP_ALIGN.CENTER)
    add_text(slide12, x + Inches(0.15), y + Inches(1.4), Inches(2.0), Inches(0.6),
             desc, size=9, color=C_MID, align=PP_ALIGN.CENTER)

add_text(slide12, Inches(0.5), Inches(4.3), Inches(12), Inches(0.3),
         "Why Docker for This Project?", size=18, color=C_DARK, bold=True)
benefits = [
    ("Isolation", "Each service has its own database\ncontainer \u2014 no schema conflicts.\nData persisted in Docker volumes.", C_PRIMARY),
    ("Reproducibility", "docker-compose up -d starts all\n5 databases with correct configs.\nSame environment everywhere.", C_ACCENT1),
    ("Database per Service", "True microservices pattern \u2014 each\nservice owns its data. No shared\ndatabase coupling between services.", C_ACCENT2),
    ("Easy Scaling", "Scale individual containers\nindependently based on load.\nAdd replicas as needed.", C_SECONDARY),
    ("Dev/Prod Parity", "Same Docker containers in dev,\nstaging, and production.\nNo 'works on my machine' issues.", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(benefits):
    x = Inches(0.3) + Inches(i * 2.6)
    y = Inches(4.7)
    add_card(slide12, x, y, Inches(2.3), Inches(1.8), title, desc, color)

add_rounded_rect(slide12, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), RGBColor(30, 41, 59))
add_text(slide12, Inches(0.7), Inches(6.73), Inches(12), Inches(0.4),
         "$ docker-compose up -d   |   5 MySQL containers   |   Persistent volumes   |   Auto-restart on failure   |   Network isolation",
         size=10, color=RGBColor(167, 243, 208))


# ══════════════════════════════════════════════════════════════
# SLIDE 13: CONCLUSION
# ══════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide13, C_WHITE)
add_header_bar(slide13, "What Makes RevConnect Unique?", "How we stand apart from Instagram, Facebook & other social platforms")
add_slide_number(slide13, 13)

comparisons = [
    ("Microservices Architecture",
     "Instagram & Facebook use monolithic or large-scale proprietary architectures. RevConnect demonstrates a "
     "clean microservices approach with 7 independent services, each with its own database \u2014 making it "
     "highly scalable, maintainable, and fault-tolerant.", C_PRIMARY),
    ("Three User Types with Business Tools",
     "Unlike Instagram/Facebook which have limited business profiles, RevConnect offers Standard, Business, "
     "and Creator accounts \u2014 each with tailored features like CTA buttons (Shop Now, Learn More), paid "
     "partnership disclosures, product tagging, and post scheduling.", C_SECONDARY),
    ("Real Analytics with Multi-Service Aggregation",
     "RevConnect's analytics dashboard aggregates data from Interaction, Connection, and Post services "
     "via Feign clients \u2014 providing real engagement metrics, follower demographics by account type, "
     "and per-post performance. Not mock data \u2014 real aggregated microservice data.", C_ACCENT1),
    ("Granular Notification Preferences",
     "Users control notifications per type (likes, comments, shares, followers) independently \u2014 not "
     "just a global on/off switch. The notification service checks user preferences before creating "
     "any notification, respecting user choice at every level.", C_ACCENT2),
    ("True Service Discovery & Gateway",
     "Uses Netflix Eureka for dynamic service registration and Spring Cloud Gateway for intelligent "
     "routing, JWT authentication, and circuit breakers \u2014 a production-grade microservices setup "
     "that Instagram/Facebook don't expose to developers.", C_ACCENT4),
    ("Open Source & Educational",
     "Built as a learning project that demonstrates full-stack microservices development with "
     "real-world patterns: database-per-service, API gateway, Feign clients, Docker containers, "
     "CI/CD pipelines, and AWS deployment \u2014 a complete portfolio project.", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(comparisons):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.3) + Inches(row * 1.9)
    add_card(slide13, x, y, Inches(5.9), Inches(1.7), title, desc, color)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: THANK YOU
# ══════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide14, C_WHITE)

# Decorative background circles
add_circle(slide14, Inches(0.3), Inches(0.3), Inches(2.5), RGBColor(238, 242, 255))
add_circle(slide14, Inches(10), Inches(5), Inches(3), RGBColor(245, 243, 255))
add_circle(slide14, Inches(-0.5), Inches(5.5), Inches(2), RGBColor(236, 253, 245))
add_circle(slide14, Inches(8.5), Inches(-0.3), Inches(1.8), RGBColor(254, 243, 199))
add_circle(slide14, Inches(11.5), Inches(2), Inches(1.5), RGBColor(254, 226, 226))

# Scattered emoji stickers around the slide (using colored shape icons instead of unicode emojis)
# Heart shapes
add_rounded_rect(slide14, Inches(0.6), Inches(0.5), Inches(0.7), Inches(0.7), C_ACCENT3)
add_text(slide14, Inches(0.6), Inches(0.52), Inches(0.7), Inches(0.7), "\u2665", size=28, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(11.8), Inches(0.6), Inches(0.7), Inches(0.7), C_PRIMARY)
add_text(slide14, Inches(11.8), Inches(0.62), Inches(0.7), Inches(0.7), "\u263a", size=28, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(0.4), Inches(3.0), Inches(0.6), Inches(0.6), C_ACCENT2)
add_text(slide14, Inches(0.4), Inches(3.02), Inches(0.6), Inches(0.6), "\u266b", size=22, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(12.3), Inches(3.3), Inches(0.6), Inches(0.6), C_ACCENT4)
add_text(slide14, Inches(12.3), Inches(3.32), Inches(0.6), Inches(0.6), "\u2606", size=22, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(0.8), Inches(5.8), Inches(0.6), Inches(0.6), C_ACCENT1)
add_text(slide14, Inches(0.8), Inches(5.82), Inches(0.6), Inches(0.6), "\u2764", size=22, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(12.0), Inches(5.5), Inches(0.6), Inches(0.6), C_SECONDARY)
add_text(slide14, Inches(12.0), Inches(5.52), Inches(0.6), Inches(0.6), "\u2605", size=22, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(2.5), Inches(0.2), Inches(0.6), Inches(0.6), C_SECONDARY)
add_text(slide14, Inches(2.5), Inches(0.22), Inches(0.6), Inches(0.6), "\u2302", size=20, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(10.3), Inches(0.2), Inches(0.6), Inches(0.6), C_ACCENT1)
add_text(slide14, Inches(10.3), Inches(0.22), Inches(0.6), Inches(0.6), "\u2709", size=20, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(5.0), Inches(0.15), Inches(0.6), Inches(0.6), C_ACCENT3)
add_text(slide14, Inches(5.0), Inches(0.17), Inches(0.6), Inches(0.6), "\u2665", size=20, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(7.8), Inches(0.15), Inches(0.6), Inches(0.6), C_PRIMARY)
add_text(slide14, Inches(7.8), Inches(0.17), Inches(0.6), Inches(0.6), "\u260e", size=20, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(3.5), Inches(6.5), Inches(0.5), Inches(0.5), C_ACCENT2)
add_text(slide14, Inches(3.5), Inches(6.52), Inches(0.5), Inches(0.5), "\u2726", size=18, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rounded_rect(slide14, Inches(9.3), Inches(6.5), Inches(0.5), Inches(0.5), C_ACCENT4)
add_text(slide14, Inches(9.3), Inches(6.52), Inches(0.5), Inches(0.5), "\u2726", size=18, color=C_WHITE, align=PP_ALIGN.CENTER)

# Main title with decorations
add_text(slide14, Inches(2), Inches(1.2), Inches(9.333), Inches(0.9),
         "\u2665  Thank You!  \u2665", size=54, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide14, Inches(2), Inches(2.1), Inches(9.333), Inches(0.5),
         "\u2606 RevConnect \u2014 A Microservices-Based Social Media Platform \u2606",
         size=20, color=C_SECONDARY, bold=True, align=PP_ALIGN.CENTER)

# Team Members section
add_rect(slide14, Inches(4.5), Inches(2.9), Inches(4.333), Inches(0.04), C_PRIMARY)
add_text(slide14, Inches(2), Inches(3.1), Inches(9.333), Inches(0.5),
         "\u2605  Team Members  \u2605", size=22, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)

team = [
    ("Pavani", C_PRIMARY),
    ("Naveen", C_ACCENT1),
    ("Ashritha", C_SECONDARY),
    ("Bala Avinash", C_ACCENT2),
    ("Pavan", C_ACCENT4),
]
for i, (name, color) in enumerate(team):
    x = Inches(1.5) + Inches(i * 2.2)
    add_rounded_rect(slide14, x, Inches(3.8), Inches(1.9), Inches(1.0), C_WHITE)
    add_rect(slide14, x, Inches(3.8), Inches(1.9), Inches(0.06), color)
    add_circle(slide14, x + Inches(0.65), Inches(3.9), Inches(0.6), color)
    add_text(slide14, x + Inches(0.65), Inches(3.95), Inches(0.6), Inches(0.5),
             name[0], size=22, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide14, x, Inches(4.5), Inches(1.9), Inches(0.3),
             name, size=13, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)

# Tech summary
add_text(slide14, Inches(3), Inches(5.1), Inches(7.333), Inches(0.8),
         "\u2726 Angular 19 + Spring Boot 3 + Docker + AWS \u2726\n"
         "\u25cb 7 Microservices  |  \u25cb 5 MySQL Databases  |  \u25cb JWT Auth",
         size=14, color=C_MID, align=PP_ALIGN.CENTER)

add_rect(slide14, Inches(5), Inches(5.95), Inches(3.333), Inches(0.04), C_PRIMARY)
add_text(slide14, Inches(2), Inches(6.1), Inches(9.333), Inches(0.5),
         "\u2605 Questions & Discussion \u2605", size=18, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(slide14, Inches(2), Inches(6.6), Inches(9.333), Inches(0.5),
         "We'd love to hear your feedback!",
         size=14, color=C_MID, align=PP_ALIGN.CENTER)
add_slide_number(slide14, 14)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = r'd:\Microservices\RevConnect_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")

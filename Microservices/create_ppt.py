from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ══════════════════════════════════════════════════════════════
# COLOR PALETTE
# ══════════════════════════════════════════════════════════════
C_PRIMARY = RGBColor(59, 130, 246)    # Blue
C_SECONDARY = RGBColor(139, 92, 246)  # Purple
C_ACCENT1 = RGBColor(16, 185, 129)   # Green
C_ACCENT2 = RGBColor(245, 158, 11)   # Amber
C_ACCENT3 = RGBColor(239, 68, 68)    # Red
C_ACCENT4 = RGBColor(6, 182, 212)    # Cyan
C_DARK = RGBColor(30, 41, 59)        # Slate 800
C_MID = RGBColor(71, 85, 105)        # Slate 600
C_LIGHT = RGBColor(241, 245, 249)    # Slate 100
C_WHITE = RGBColor(255, 255, 255)
C_BG_WARM = RGBColor(255, 251, 235)  # Warm white
C_GRADIENT_START = RGBColor(238, 242, 255)  # Indigo 50

# ══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════
def add_bg(slide, color=C_WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=C_DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body, icon_color, title_color=C_DARK):
    card = add_rounded_rect(slide, left, top, width, height, C_WHITE)
    card.shadow.inherit = False
    # Icon strip at top
    add_rect(slide, left, top, width, Inches(0.06), icon_color)
    # Title
    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.35),
             title, size=13, color=title_color, bold=True)
    # Body
    add_text(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
             body, size=10, color=C_MID)
    return card

def add_slide_number(slide, num, total=10):
    add_text(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
             f"{num}/{total}", size=10, color=C_MID, align=PP_ALIGN.RIGHT)

def add_header_bar(slide, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), C_PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
             title, size=28, color=C_WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4),
                 subtitle, size=14, color=RGBColor(191, 219, 254))
    # Decorative circles
    add_circle(slide, Inches(11.5), Inches(-0.3), Inches(0.8), RGBColor(96, 165, 250))
    add_circle(slide, Inches(12.2), Inches(0.3), Inches(0.6), C_SECONDARY)

# ══════════════════════════════════════════════════════════════
# SLIDE 1: ABSTRACT (TITLE PAGE)
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide1, C_WHITE)

# Large gradient-like background block
add_rect(slide1, Inches(0), Inches(0), Inches(5.5), prs.slide_height, RGBColor(238, 242, 255))
add_rect(slide1, Inches(0), Inches(0), Inches(0.08), prs.slide_height, C_PRIMARY)

# Decorative circles
add_circle(slide1, Inches(10), Inches(0.5), Inches(2.5), RGBColor(238, 242, 255))
add_circle(slide1, Inches(11), Inches(5), Inches(2), RGBColor(245, 243, 255))
add_circle(slide1, Inches(4.5), Inches(5.5), Inches(1.5), RGBColor(236, 253, 245))

# Title
add_text(slide1, Inches(0.8), Inches(0.8), Inches(4), Inches(0.5),
         "RevConnect", size=42, color=C_PRIMARY, bold=True)
add_text(slide1, Inches(0.8), Inches(1.4), Inches(4), Inches(0.5),
         "Social Media Platform", size=22, color=C_SECONDARY, bold=True)

# Abstract box
add_text(slide1, Inches(0.8), Inches(2.2), Inches(4.2), Inches(0.3),
         "ABSTRACT", size=14, color=C_PRIMARY, bold=True)

abstract_text = (
    "RevConnect is a full-stack professional social media platform built using a microservices architecture "
    "with Angular 19 frontend and Spring Boot 3 backend. The system comprises 7 independent microservices "
    "(User, Post, Interaction, Connection, Notification, Feed, API Gateway) each with its own MySQL database, "
    "communicating via REST APIs and Eureka service discovery.\n\n"
    "Key features include user registration with email verification, JWT-based authentication, real-time "
    "personalized feeds, post creation with media uploads (images/videos), social interactions (likes, comments, "
    "shares), connection management (follow/unfollow), direct messaging with WebSocket support, push/email "
    "notifications with granular preferences, user analytics dashboards, stories, bookmarks, explore/search, "
    "and business tools (CTA buttons, paid partnerships, product tagging, scheduled posts).\n\n"
    "The platform supports three user types: Standard, Business, and Creator — each with tailored features. "
    "Privacy controls allow users to set profiles as Public or Private. The architecture ensures scalability, "
    "fault tolerance via circuit breakers, and clean separation of concerns."
)
add_text(slide1, Inches(0.8), Inches(2.6), Inches(4.2), Inches(4.5),
         abstract_text, size=11, color=C_MID)

# Right side — Tech stack cards
add_text(slide1, Inches(6), Inches(0.8), Inches(6), Inches(0.4),
         "Technology Stack", size=20, color=C_DARK, bold=True)

techs = [
    ("Frontend", "Angular 19, TypeScript, RxJS, TailwindCSS, Font Awesome", C_PRIMARY),
    ("Backend", "Spring Boot 3, Java 17, Spring Cloud, Eureka, Feign, JWT", C_SECONDARY),
    ("Database", "MySQL 8 (5 Docker containers), JPA/Hibernate, Spring Data", C_ACCENT1),
    ("DevOps", "Docker Compose, API Gateway, Circuit Breakers, Eureka Discovery", C_ACCENT2),
    ("Features", "Real-time Feed, Analytics, Stories, Messages, Notifications, Business Tools", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(techs):
    y = Inches(1.4) + Inches(i * 1.1)
    add_card(slide1, Inches(6), y, Inches(6.8), Inches(0.9), title, desc, color)

add_slide_number(slide1, 1)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, C_WHITE)
add_header_bar(slide2, "System Architecture", "Microservices Architecture with Service Discovery & API Gateway")
add_slide_number(slide2, 2)

# Client Layer
add_rounded_rect(slide2, Inches(5.2), Inches(1.3), Inches(3), Inches(0.7), RGBColor(219, 234, 254))
add_text(slide2, Inches(5.2), Inches(1.35), Inches(3), Inches(0.35),
         "Angular 19 Frontend", size=14, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide2, Inches(5.2), Inches(1.65), Inches(3), Inches(0.3),
         "localhost:4200", size=10, color=C_MID, align=PP_ALIGN.CENTER)

# Arrow down
add_arrow(slide2, Inches(6.4), Inches(2.1), Inches(0.6), Inches(0.3), C_MID)

# API Gateway
add_rounded_rect(slide2, Inches(4.5), Inches(2.5), Inches(4.4), Inches(0.8), RGBColor(254, 243, 199))
add_text(slide2, Inches(4.5), Inches(2.55), Inches(4.4), Inches(0.4),
         "API Gateway (Spring Cloud Gateway)", size=14, color=C_ACCENT2, bold=True, align=PP_ALIGN.CENTER)
add_text(slide2, Inches(4.5), Inches(2.9), Inches(4.4), Inches(0.3),
         "JWT Auth | Route | Load Balance | Circuit Breaker — :8080", size=10, color=C_MID, align=PP_ALIGN.CENTER)

# Eureka
add_rounded_rect(slide2, Inches(10), Inches(2.5), Inches(2.8), Inches(0.8), RGBColor(237, 233, 254))
add_text(slide2, Inches(10), Inches(2.55), Inches(2.8), Inches(0.4),
         "Eureka Server", size=14, color=C_SECONDARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide2, Inches(10), Inches(2.9), Inches(2.8), Inches(0.3),
         "Service Discovery — :8761", size=10, color=C_MID, align=PP_ALIGN.CENTER)

# Services row
services = [
    ("User\nService", ":8081", C_PRIMARY, "Auth, Profile\nSettings, Search"),
    ("Post\nService", ":8082", C_ACCENT1, "CRUD, Schedule\nMedia, CTA"),
    ("Feed\nService", ":8083", C_ACCENT4, "Personalized\nFeed Aggregation"),
    ("Interaction\nService", ":8084", C_ACCENT2, "Likes, Comments\nShares, Views"),
    ("Connection\nService", ":8085", C_SECONDARY, "Follow/Unfollow\nStats, Network"),
    ("Notification\nService", ":8086", C_ACCENT3, "Push, Email\nPreferences"),
]

for i, (name, port, color, desc) in enumerate(services):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(4.0)
    card = add_rounded_rect(slide2, x, y, Inches(1.95), Inches(1.8), C_WHITE)
    # Top color bar
    add_rect(slide2, x, y, Inches(1.95), Inches(0.08), color)
    add_text(slide2, x, y + Inches(0.15), Inches(1.95), Inches(0.5),
             name, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide2, x, y + Inches(0.65), Inches(1.95), Inches(0.25),
             port, size=9, color=C_MID, align=PP_ALIGN.CENTER)
    add_text(slide2, x, y + Inches(0.9), Inches(1.95), Inches(0.7),
             desc, size=9, color=C_MID, align=PP_ALIGN.CENTER)

# Database row
dbs = [
    ("user_db", ":3307", C_PRIMARY),
    ("post_db", ":3308", C_ACCENT1),
    ("—", "", C_ACCENT4),
    ("interaction_db", ":3309", C_ACCENT2),
    ("connection_db", ":3310", C_SECONDARY),
    ("notification_db", ":3311", C_ACCENT3),
]
for i, (name, port, color) in enumerate(dbs):
    if name == "—":
        continue
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(6.1)
    shape = add_rounded_rect(slide2, x, y, Inches(1.95), Inches(0.7), RGBColor(241, 245, 249))
    add_text(slide2, x, y + Inches(0.05), Inches(1.95), Inches(0.35),
             f"MySQL {name}", size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide2, x, y + Inches(0.35), Inches(1.95), Inches(0.25),
             f"Docker {port}", size=8, color=C_MID, align=PP_ALIGN.CENTER)

# Label
add_text(slide2, Inches(0.3), Inches(6.9), Inches(12), Inches(0.3),
         "Each microservice has its own isolated MySQL database running in Docker containers — Database per Service pattern",
         size=10, color=C_MID, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: LOGIN & REGISTER
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, C_WHITE)
add_header_bar(slide3, "Authentication Flow", "Login, Register & Forgot Password")
add_slide_number(slide3, 3)

# LOGIN FLOW
add_text(slide3, Inches(0.5), Inches(1.3), Inches(5), Inches(0.4),
         "Login Flow", size=18, color=C_PRIMARY, bold=True)

login_steps = [
    ("1", "User enters\nemail & password", C_PRIMARY),
    ("2", "Angular sends\nPOST /api/auth/login", C_ACCENT4),
    ("3", "Gateway routes\nto User Service", C_ACCENT2),
    ("4", "BCrypt verifies\npassword hash", C_ACCENT3),
    ("5", "JWT token\ngenerated & returned", C_ACCENT1),
]
for i, (num, text, color) in enumerate(login_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_circle(slide3, x, Inches(1.9), Inches(0.5), color)
    add_text(slide3, x, Inches(1.93), Inches(0.5), Inches(0.5),
             num, size=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, x - Inches(0.3), Inches(2.5), Inches(1.1), Inches(0.6),
             text, size=10, color=C_DARK, align=PP_ALIGN.CENTER)
    if i < len(login_steps) - 1:
        add_arrow(slide3, x + Inches(0.6), Inches(2.05), Inches(1.6), Inches(0.2), RGBColor(203, 213, 225))

add_text(slide3, Inches(0.5), Inches(3.2), Inches(12), Inches(0.5),
         "Dependencies: BCryptPasswordEncoder, JwtTokenProvider (HS256, 256-bit secret), Spring Security FilterChain",
         size=10, color=C_MID)

# REGISTER FLOW
add_text(slide3, Inches(0.5), Inches(3.8), Inches(5), Inches(0.4),
         "Register Flow", size=18, color=C_ACCENT1, bold=True)

reg_steps = [
    ("1", "User fills form:\nname, email, password", C_ACCENT1),
    ("2", "POST /api/auth/\nregister", C_ACCENT4),
    ("3", "Duplicate check\n+ BCrypt hash", C_ACCENT2),
    ("4", "Verification email\nsent (code in DB)", C_ACCENT3),
    ("5", "User clicks link\naccount activated", C_PRIMARY),
]
for i, (num, text, color) in enumerate(reg_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_circle(slide3, x, Inches(4.4), Inches(0.5), color)
    add_text(slide3, x, Inches(4.43), Inches(0.5), Inches(0.5),
             num, size=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, x - Inches(0.3), Inches(5.0), Inches(1.1), Inches(0.6),
             text, size=10, color=C_DARK, align=PP_ALIGN.CENTER)
    if i < len(reg_steps) - 1:
        add_arrow(slide3, x + Inches(0.6), Inches(4.55), Inches(1.6), Inches(0.2), RGBColor(203, 213, 225))

# FORGOT PASSWORD
add_text(slide3, Inches(0.5), Inches(5.8), Inches(5), Inches(0.4),
         "Forgot Password Flow", size=18, color=C_ACCENT3, bold=True)
add_text(slide3, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
         "User enters email → POST /api/auth/forgot-password → User Service generates reset token (UUID) → "
         "Email sent with reset link → User clicks link with token → Enters new password → "
         "POST /api/auth/reset-password with token → Password re-hashed with BCrypt → Success",
         size=11, color=C_MID)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: HOME PAGE & FEED
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, C_WHITE)
add_header_bar(slide4, "Home Page & Personalized Feed", "How posts are loaded and displayed")
add_slide_number(slide4, 4)

# Feed flow
flow_items = [
    ("User opens /feed", "FeedPage component\ninitializes", C_PRIMARY),
    ("Load profile", "GET /api/users/me\nvia API Gateway + JWT", C_ACCENT4),
    ("Fetch feed", "GET /api/posts/feed/\npersonalized", C_ACCENT1),
    ("Feed Service", "Gets following IDs from\nConnection Service (Feign)", C_ACCENT2),
    ("Post Service", "Queries posts by\nfollowing user IDs", C_SECONDARY),
    ("Enrich posts", "Fetches author info\nfrom User Service", C_ACCENT3),
]

for i, (title, desc, color) in enumerate(flow_items):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(1.4)
    add_rounded_rect(slide4, x, y, Inches(1.95), Inches(1.5), C_WHITE)
    add_rect(slide4, x, y, Inches(1.95), Inches(0.06), color)
    add_text(slide4, x + Inches(0.1), y + Inches(0.15), Inches(1.75), Inches(0.3),
             title, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide4, x + Inches(0.1), y + Inches(0.5), Inches(1.75), Inches(0.8),
             desc, size=10, color=C_MID, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_arrow(slide4, x + Inches(1.95), y + Inches(0.6), Inches(0.2), Inches(0.15), RGBColor(203, 213, 225))

# Additional details
add_text(slide4, Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
         "Feed Page Components & Features", size=16, color=C_DARK, bold=True)

cards_data = [
    ("Stories Feed", "Horizontal scrollable story cards at top.\nUsers can create and view 24h stories.", C_PRIMARY),
    ("Create Post Box", "Text input + Media upload + Schedule +\nCTA buttons. Supports IMAGE/VIDEO/TEXT types.", C_ACCENT1),
    ("Live Analytics Sidebar", "Shows total views, likes, shares, comments\nin real-time for Business/Creator users.", C_ACCENT2),
    ("Trending Topics", "Fetched from search service. Shows\npopular hashtags and search terms.", C_SECONDARY),
    ("Post Cards", "Each post shows author info, content, media,\nlike/comment/share buttons, timestamps.", C_ACCENT3),
    ("Infinite Scroll", "Paginated feed (page=0, size=10).\nLoads more posts as user scrolls.", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(cards_data):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(3.7) + Inches(row * 1.7)
    add_card(slide4, x, y, Inches(3.8), Inches(1.5), title, desc, color)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: POSTS — CREATE, SCHEDULE, CTA
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, C_WHITE)
add_header_bar(slide5, "Post Creation & Business Tools", "Create, Schedule, CTA, Partnerships & Product Tags")
add_slide_number(slide5, 5)

# Normal post flow
add_text(slide5, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
         "Normal Post Creation", size=16, color=C_PRIMARY, bold=True)

post_steps = [
    ("Write content\n+ attach media", C_PRIMARY),
    ("Upload media to\nPOST /api/media/upload", C_ACCENT4),
    ("POST /api/posts\nwith content + mediaUrls", C_ACCENT1),
    ("PostService saves\nto MySQL (post_db)", C_ACCENT2),
    ("Post appears\nin feed", C_SECONDARY),
]
for i, (text, color) in enumerate(post_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_rounded_rect(slide5, x, Inches(1.7), Inches(2.1), Inches(0.8), color)
    add_text(slide5, x + Inches(0.1), Inches(1.75), Inches(1.9), Inches(0.7),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(post_steps) - 1:
        add_arrow(slide5, x + Inches(2.1), Inches(1.95), Inches(0.4), Inches(0.15), RGBColor(203, 213, 225))

# Scheduled post flow
add_text(slide5, Inches(0.5), Inches(2.8), Inches(5), Inches(0.3),
         "Scheduled Post with CTA", size=16, color=C_ACCENT2, bold=True)

sched_steps = [
    ("User selects date,\ntime, CTA label & URL", C_ACCENT2),
    ("Frontend builds ISO\ndatetime + request body", C_ACCENT4),
    ("POST /api/posts/\nschedule", C_ACCENT1),
    ("Post saved with\nisPublished=false", C_ACCENT3),
    ("Scheduler publishes\nat scheduled time", C_SECONDARY),
]
for i, (text, color) in enumerate(sched_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_rounded_rect(slide5, x, Inches(3.2), Inches(2.1), Inches(0.8), color)
    add_text(slide5, x + Inches(0.1), Inches(3.25), Inches(1.9), Inches(0.7),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(sched_steps) - 1:
        add_arrow(slide5, x + Inches(2.1), Inches(3.45), Inches(0.4), Inches(0.15), RGBColor(203, 213, 225))

# Business vs Creator
add_text(slide5, Inches(0.5), Inches(4.3), Inches(5), Inches(0.3),
         "User Type Differences", size=16, color=C_DARK, bold=True)

add_card(slide5, Inches(0.5), Inches(4.7), Inches(5.8), Inches(2.4),
         "BUSINESS Account",
         "• Post Categories: Standard, Announcement, Product Update\n"
         "• Call To Action: Shop Now, Learn More, Sign Up, Contact Us,\n  Download, Subscribe, Book Now\n"
         "• CTA URL: Link to external page\n"
         "• No paid partnership or product tagging",
         C_PRIMARY)

add_card(slide5, Inches(6.8), Inches(4.7), Inches(6), Inches(2.4),
         "CREATOR Account",
         "• All Business features PLUS:\n"
         "• Paid Partnership Disclosure: Mark as sponsored,\n  enter brand partner name (e.g. Nike, Apple)\n"
         "• Product Tagging: Comma-separated tags\n  (e.g. Shoes, Watch, Sunglasses)\n"
         "• Partnership info stored in partnerName & productTags fields",
         C_SECONDARY)

# ══════════════════════════════════════════════════════════════
# SLIDE 6: LIKES, COMMENTS, SHARES
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, C_WHITE)
add_header_bar(slide6, "Social Interactions", "Likes, Comments & Shares — Interaction Service")
add_slide_number(slide6, 6)

# Likes
add_text(slide6, Inches(0.5), Inches(1.3), Inches(4), Inches(0.3),
         "Like Flow", size=16, color=C_ACCENT3, bold=True)
like_flow = (
    "User clicks ❤ → Frontend calls POST /api/posts/{id}/like → "
    "API Gateway routes to Interaction Service → "
    "Checks if already liked (unique constraint: user_id + post_id) → "
    "If not liked: saves Like entity, sends notification to post author via Notification Service (Feign) → "
    "If already liked: returns 'Already liked'. "
    "Unlike: DELETE /api/posts/{id}/like removes the Like record."
)
add_text(slide6, Inches(0.5), Inches(1.7), Inches(12), Inches(1.0),
         like_flow, size=11, color=C_MID)

# Comments
add_text(slide6, Inches(0.5), Inches(2.8), Inches(4), Inches(0.3),
         "Comment Flow", size=16, color=C_PRIMARY, bold=True)
comment_flow = (
    "User types comment → POST /api/posts/{id}/comments with {content, parentId?} → "
    "Interaction Service saves Comment entity (supports nested replies via parentId) → "
    "Notification sent to post author → Comments fetched with pagination: "
    "GET /api/posts/{id}/comments?page=0&size=10 → "
    "Each comment enriched with author name/avatar from User Service."
)
add_text(slide6, Inches(0.5), Inches(3.2), Inches(12), Inches(1.0),
         comment_flow, size=11, color=C_MID)

# Shares
add_text(slide6, Inches(0.5), Inches(4.3), Inches(4), Inches(0.3),
         "Share Flow", size=16, color=C_ACCENT1, bold=True)
share_flow = (
    "User clicks Share → POST /api/posts/{id}/share with optional {comment} → "
    "Creates Share entity with REPOST type → Notification sent to original post author → "
    "Share count visible in post cards and analytics. "
    "Shares tracked per user in shares table."
)
add_text(slide6, Inches(0.5), Inches(4.7), Inches(12), Inches(1.0),
         share_flow, size=11, color=C_MID)

# View tracking
add_text(slide6, Inches(0.5), Inches(5.6), Inches(4), Inches(0.3),
         "View Tracking", size=16, color=C_ACCENT4, bold=True)
view_flow = (
    "When a user's feed loads, each post records a view: POST /api/posts/{id}/view → "
    "PostView entity with unique constraint (post_id + user_id) ensures each user counted only once per post → "
    "View count used in analytics dashboard. No duplicate counting on page refresh."
)
add_text(slide6, Inches(0.5), Inches(6.0), Inches(12), Inches(1.0),
         view_flow, size=11, color=C_MID)

# Dependencies box
add_rounded_rect(slide6, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), RGBColor(241, 245, 249))
add_text(slide6, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "Dependencies: LikeRepository, CommentRepository, ShareRepository, PostViewRepository, NotificationServiceClient (Feign), UserServiceClient (Feign)",
         size=10, color=C_MID)

# ══════════════════════════════════════════════════════════════
# SLIDE 7: MESSAGES
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, C_WHITE)
add_header_bar(slide7, "Direct Messages", "Real-time messaging between connected users")
add_slide_number(slide7, 7)

msg_cards = [
    ("Conversation List", "GET /api/messages/conversations\nReturns list of conversations with last message preview,\ntimestamp, unread count. Sorted by newest first.", C_PRIMARY),
    ("Load Messages", "GET /api/messages/conversation/{recipientId}\nFetches paginated messages between two users.\nMessages shown in chronological order.", C_ACCENT4),
    ("Send Message", "POST /api/messages with {recipientId, content, type}\nSupports TEXT, IMAGE, VOICE message types.\nNotification sent to recipient.", C_ACCENT1),
    ("Real-time Updates", "WebSocket connection via STOMP protocol.\nNew messages pushed instantly to recipients.\nNo polling needed — server pushes updates.", C_ACCENT2),
    ("Message Features", "Read receipts, typing indicators, message deletion.\nVoice messages recorded in browser and uploaded.\nImage messages use media upload service.", C_SECONDARY),
    ("Architecture", "Messages stored in User Service database.\nFrontend uses RxJS observables for reactive updates.\nConversations sorted by lastMessageTime descending.", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(msg_cards):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.4) + Inches(row * 2.8)
    add_card(slide7, x, y, Inches(3.8), Inches(2.5), title, desc, color)

# ══════════════════════════════════════════════════════════════
# SLIDE 8: EXPLORE & CONNECTIONS
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, C_WHITE)
add_header_bar(slide8, "Explore & Connections", "Search, Follow/Unfollow & Network Management")
add_slide_number(slide8, 8)

# Explore
add_text(slide8, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
         "Explore / Search", size=16, color=C_PRIMARY, bold=True)

search_steps = [
    ("User types query\nin search bar", C_PRIMARY),
    ("GET /api/search/all?\nquery=... (debounced)", C_ACCENT4),
    ("SearchController calls\nUserService + PostService", C_ACCENT1),
    ("Users filtered: PRIVATE\naccounts excluded", C_ACCENT3),
    ("Results shown in\nAll / Users / Posts tabs", C_SECONDARY),
]
for i, (text, color) in enumerate(search_steps):
    x = Inches(0.3) + Inches(i * 2.55)
    add_rounded_rect(slide8, x, Inches(1.7), Inches(2.2), Inches(0.9), color)
    add_text(slide8, x + Inches(0.1), Inches(1.75), Inches(2.0), Inches(0.8),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(search_steps) - 1:
        add_arrow(slide8, x + Inches(2.2), Inches(2.0), Inches(0.35), Inches(0.15), RGBColor(203, 213, 225))

# Connections
add_text(slide8, Inches(0.5), Inches(2.9), Inches(5), Inches(0.3),
         "Follow / Unfollow Flow", size=16, color=C_ACCENT1, bold=True)

conn_steps = [
    ("User clicks Follow\non profile card", C_ACCENT1),
    ("POST /api/users/\n{id}/follow", C_ACCENT4),
    ("Connection Service\ncreates Follow record", C_ACCENT2),
    ("FOLLOW notification\nsent to target user", C_ACCENT3),
    ("Follower/Following\ncounts updated", C_PRIMARY),
]
for i, (text, color) in enumerate(conn_steps):
    x = Inches(0.3) + Inches(i * 2.55)
    add_rounded_rect(slide8, x, Inches(3.3), Inches(2.2), Inches(0.9), color)
    add_text(slide8, x + Inches(0.1), Inches(3.35), Inches(2.0), Inches(0.8),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(conn_steps) - 1:
        add_arrow(slide8, x + Inches(2.2), Inches(3.6), Inches(0.35), Inches(0.15), RGBColor(203, 213, 225))

# Privacy & Suggestions
add_text(slide8, Inches(0.5), Inches(4.5), Inches(5), Inches(0.3),
         "Privacy & Suggestions", size=16, color=C_DARK, bold=True)

priv_cards = [
    ("Privacy Controls",
     "Users can set profile to PUBLIC or PRIVATE.\n"
     "PRIVATE users are excluded from:\n"
     "• Search results (searchUsers query filter)\n"
     "• Suggested users (findSuggestions query filter)\n"
     "• Explore page discovery\n"
     "Privacy updated via PATCH /api/users/me/privacy",
     C_ACCENT3),
    ("Suggested Users",
     "GET /api/users/suggested returns non-private,\n"
     "email-verified users excluding current user.\n"
     "Shown in Explore page sidebar.\n"
     "Algorithm: Basic exclusion filter.\n"
     "Future: ML-based recommendation engine.",
     C_PRIMARY),
    ("Connection Stats",
     "GET /api/users/{id}/connection-stats returns:\n"
     "• followersCount — users following you\n"
     "• followingCount — users you follow\n"
     "Used in profile page and analytics.\n"
     "Stored in Connection Service (connection_db).",
     C_ACCENT1),
]
for i, (title, desc, color) in enumerate(priv_cards):
    x = Inches(0.5) + Inches(i * 4.2)
    add_card(slide8, x, Inches(4.9), Inches(3.8), Inches(2.4), title, desc, color)

# ══════════════════════════════════════════════════════════════
# SLIDE 9: PROFILE & NOTIFICATIONS
# ══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, C_WHITE)
add_header_bar(slide9, "Profile & Notifications", "User profile management and notification system")
add_slide_number(slide9, 9)

# Profile
add_text(slide9, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
         "Profile Page Features", size=16, color=C_PRIMARY, bold=True)

profile_cards = [
    ("Profile Display", "Cover photo, avatar, name, bio, location,\nwebsite, user type badge, follower stats.\nGET /api/users/{id} + connection-stats", C_PRIMARY),
    ("Edit Profile", "PUT /api/users/me with ProfileUpdateRequest.\nUpdates name, bio, location, avatar,\ncover photo, business info, social links.", C_ACCENT1),
    ("User Posts Tab", "GET /api/posts/user/{id} — paginated.\nShows all posts by the user with\nlike/comment/share counts.", C_ACCENT2),
    ("Bookmarks", "GET /api/interactions/bookmarks.\nSaved posts for later reading.\nToggle via POST/DELETE bookmark endpoint.", C_SECONDARY),
]
for i, (title, desc, color) in enumerate(profile_cards):
    x = Inches(0.5) + Inches(i * 3.15)
    add_card(slide9, x, Inches(1.7), Inches(2.9), Inches(1.8), title, desc, color)

# Notifications
add_text(slide9, Inches(0.5), Inches(3.8), Inches(5), Inches(0.3),
         "Notification System", size=16, color=C_ACCENT3, bold=True)

notif_steps = [
    ("Action triggers\nnotification", C_PRIMARY),
    ("Interaction Service\ncalls Notification API", C_ACCENT4),
    ("Check user prefs\n(granular per-type)", C_ACCENT2),
    ("If allowed: save\nNotification entity", C_ACCENT1),
    ("Frontend polls\nunread count", C_ACCENT3),
]
for i, (text, color) in enumerate(notif_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    add_rounded_rect(slide9, x, Inches(4.2), Inches(2.1), Inches(0.8), color)
    add_text(slide9, x + Inches(0.1), Inches(4.25), Inches(1.9), Inches(0.7),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(notif_steps) - 1:
        add_arrow(slide9, x + Inches(2.1), Inches(4.45), Inches(0.4), Inches(0.15), RGBColor(203, 213, 225))

# Notification preferences
add_text(slide9, Inches(0.5), Inches(5.3), Inches(12), Inches(0.3),
         "Granular Notification Preferences", size=14, color=C_DARK, bold=True)

pref_items = [
    ("Push Notifications", "Global on/off toggle", C_SECONDARY),
    ("Likes", "notifyLike", C_ACCENT3),
    ("Comments", "notifyComment", C_PRIMARY),
    ("New Followers", "notifyNewFollower", C_ACCENT1),
    ("Shares", "notifyShare", C_ACCENT4),
    ("Connections", "notifyConnectionRequest", C_ACCENT2),
]
for i, (title, pref, color) in enumerate(pref_items):
    x = Inches(0.5) + Inches(i * 2.1)
    add_rounded_rect(slide9, x, Inches(5.7), Inches(1.9), Inches(0.8), C_WHITE)
    add_rect(slide9, x, Inches(5.7), Inches(1.9), Inches(0.05), color)
    add_text(slide9, x + Inches(0.1), Inches(5.8), Inches(1.7), Inches(0.3),
             title, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide9, x + Inches(0.1), Inches(6.1), Inches(1.7), Inches(0.3),
             pref, size=9, color=C_MID, align=PP_ALIGN.CENTER)

add_text(slide9, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
         "Types: LIKE, COMMENT, SHARE, FOLLOW, MENTION, SYSTEM | Checked in NotificationController before saving | Stored in UserSettings entity",
         size=10, color=C_MID)

# ══════════════════════════════════════════════════════════════
# SLIDE 10: ANALYTICS
# ══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, C_WHITE)
add_header_bar(slide10, "Analytics Dashboard", "Data aggregation from multiple services")
add_slide_number(slide10, 10)

# How analytics works
add_text(slide10, Inches(0.5), Inches(1.3), Inches(12), Inches(0.3),
         "How Analytics Collects Data", size=16, color=C_PRIMARY, bold=True)

analytics_flow = [
    ("User opens\nAnalytics page", C_PRIMARY),
    ("GET /api/analytics/\noverview", C_ACCENT4),
    ("Fetch follower count\nfrom Connection Service", C_ACCENT1),
    ("Fetch user's post IDs\nfrom Post Service", C_ACCENT2),
    ("Count likes, comments,\nshares, views per post", C_ACCENT3),
    ("Aggregate &\nreturn totals", C_SECONDARY),
]
for i, (text, color) in enumerate(analytics_flow):
    x = Inches(0.3) + Inches(i * 2.15)
    add_rounded_rect(slide10, x, Inches(1.7), Inches(1.95), Inches(0.9), color)
    add_text(slide10, x + Inches(0.1), Inches(1.75), Inches(1.75), Inches(0.8),
             text, size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(analytics_flow) - 1:
        add_arrow(slide10, x + Inches(1.95), Inches(2.0), Inches(0.2), Inches(0.15), RGBColor(203, 213, 225))

# Metrics cards
add_text(slide10, Inches(0.5), Inches(2.9), Inches(12), Inches(0.3),
         "Analytics Endpoints & Metrics", size=16, color=C_DARK, bold=True)

metrics = [
    ("Overview", "/api/analytics/overview\n\nReturns: totalViews, totalLikes,\ntotalComments, totalShares,\ntotalFollowers, totalPosts", C_PRIMARY),
    ("Post Performance", "/api/analytics/posts\n\nPer-post breakdown: likes, comments,\nshares, views for each post.\nTop 5 posts by engagement.", C_ACCENT1),
    ("Follower Growth", "/api/analytics/followers\n\nFollower count over time.\nNew followers per day/week.\nGrowth trend data.", C_ACCENT2),
    ("Engagement Rate", "/api/analytics/engagement\n\nTotal engagement across all posts.\nEngagement = likes + comments + shares.\nRate = engagement / views.", C_ACCENT3),
    ("Audience", "/api/analytics/audience\n\nDemographic breakdown of followers.\nIndustry, location distribution.\nActive hours analysis.", C_SECONDARY),
    ("Profile Views", "/api/analytics/profile-views\n\nTrack who views your profile.\nReal view counting with\nPostView entity (unique per user).", C_ACCENT4),
]
for i, (title, desc, color) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(3.3) + Inches(row * 2.1)
    add_card(slide10, x, y, Inches(3.8), Inches(1.9), title, desc, color)

# Footer
add_rounded_rect(slide10, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35), RGBColor(241, 245, 249))
add_text(slide10, Inches(0.7), Inches(7.03), Inches(12), Inches(0.3),
         "Dependencies: AnalyticsController (Interaction Service) → PostServiceClient (Feign) → ConnectionServiceClient (Feign) → LikeRepo, CommentRepo, ShareRepo, PostViewRepo",
         size=9, color=C_MID)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = r'd:\Microservices\RevConnect_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
